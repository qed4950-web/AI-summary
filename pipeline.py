# pipeline.py  (Step2: 추출 + 학습)
import importlib
import math
import os, re, sys, time, threading, platform
from datetime import datetime
import io
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from dataclasses import dataclass
from typing import Optional, Dict, Any, List, Tuple, Union, Set

import numpy as np

# ---- 선택 의존성(있으면 사용) ----
try:
    import pandas as pd
except Exception:
    pd = None
PARQUET_ENGINE: Optional[str] = None
if pd is not None:
    for candidate in ("fastparquet", "pyarrow"):
        try:
            importlib.import_module(candidate)
            PARQUET_ENGINE = candidate
            break
        except ImportError:
            continue
try:
    from deep_translator import GoogleTranslator
except Exception:
    GoogleTranslator = None
try:
    import docx
except Exception:
    docx = None
try:
    import pptx
except Exception:
    pptx = None
try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None
try:
    import win32com.client
except Exception:
    win32com = None
try:
    import pythoncom
except Exception:
    pythoncom = None
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None
try:
    import joblib
except Exception:
    joblib = None
try:
    import pdfplumber
except Exception:
    pdfplumber = None
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.decomposition import TruncatedSVD
    from sklearn.cluster import MiniBatchKMeans
    from sklearn.pipeline import Pipeline
    from sklearn import __version__ as sklearn_version
except Exception:
    TfidfVectorizer = TruncatedSVD = MiniBatchKMeans = Pipeline = None
    sklearn_version = "0"

try:
    from tqdm import tqdm
except Exception:
    tqdm = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None
try:
    import olefile
except Exception:
    olefile = None
try:
    import pyhwp
except Exception:
    pyhwp = None


# =========================
# 콘솔 진행도 유틸
# =========================
class Spinner:
    FRAMES = ["|", "/", "-", "\\"]
    def __init__(self, prefix="", interval=0.12):
        self.prefix = prefix
        self.interval = interval
        self._stop = threading.Event()
        self._t = None
        self._i = 0
    def start(self):
        if self._t: return
        def _run():
            while not self._stop.wait(self.interval):
                frame = self.FRAMES[self._i % len(self.FRAMES)]
                self._i += 1
                sys.stdout.write(f"\r{self.prefix} {frame} ")
                sys.stdout.flush()
        self._t = threading.Thread(target=_run, daemon=True)
        self._t.start()
    def stop(self, clear=True):
        if not self._t: return
        self._stop.set()
        self._t.join()
        if clear:
            sys.stdout.write("\r" + " " * 80 + "\r")
            sys.stdout.flush()

class ProgressLine:
    def __init__(self, total:int, label:str, update_every:int=10):
        self.total = max(1, total)
        self.label = label
        self.update_every = max(1, update_every)
        self.start = time.time()
        self.n = 0
    def update(self, k:int=1):
        self.n += k
        if (self.n % self.update_every) != 0 and self.n < self.total:
            return
        pct = min(100.0, self.n / self.total * 100.0)
        elapsed = time.time() - self.start
        rate = self.n/elapsed if elapsed>0 else 0
        remain = (self.total - self.n)/rate if rate>0 else 0
        sys.stdout.write(
            f"\r[{pct:5.1f}%] {self.label}  {self.n:,}/{self.total:,}  "
            f"{rate:,.1f}/s  elapsed={self._fmt(elapsed)}  ETA={self._fmt(remain)}   "
        )
        sys.stdout.flush()
    def close(self):
        self.n = self.total
        self.update(0)
        sys.stdout.write("\n"); sys.stdout.flush()
    @staticmethod
    def _fmt(s: float)->str:
        if s==float("inf"): return "∞"
        m, sec = divmod(int(s), 60); h, m = divmod(m, 60)
        return f"{h:d}:{m:02d}:{sec:02d}" if h else f"{m:02d}:{sec:02d}"


# =========================
# 텍스트 클린
# =========================
class TextCleaner:
    _multi = re.compile(r"\s+")
    @classmethod
    def clean(cls, s:str)->str:
        if not s: return ""
        s = "".join(ch if ch.isprintable() or ch in "\t\n\r" else " " for ch in s)
        s = s.replace("\x00"," ")
        return cls._multi.sub(" ", s).strip()

TOKEN_PATTERN = r'(?u)(?:[가-힣]{1,}|[A-Za-z0-9]{2,})'

# 고정된 SVD 차원 수. Index/모델 불일치를 막기 위해 한곳에서 정의한다.
DEFAULT_N_COMPONENTS = 128
MODEL_TEXT_COLUMN = "text_model"
_META_SPLIT_RE = re.compile(r"[^0-9A-Za-z가-힣]+")
DEFAULT_EMBED_MODEL = "sentence-transformers/paraphrase-multilingual-MiniLM-L12-v2"
MODEL_TYPE_SENTENCE_TRANSFORMER = "sentence-transformer"

DEFAULT_CHUNK_MIN_TOKENS = 200
DEFAULT_CHUNK_MAX_TOKENS = 500

_TOKEN_REGEX = re.compile(TOKEN_PATTERN)

_EN_STOPWORDS: Set[str] = {
    "the", "and", "for", "that", "with", "from", "this", "have", "been", "were",
    "into", "about", "after", "before", "while", "shall", "could", "would", "there",
    "their", "which", "should", "among", "within", "between", "through", "without",
    "because", "against", "during", "under", "over", "where", "when", "whose", "them",
    "they", "these", "those", "ours", "your", "yours", "ourselves", "yourself",
    "yourselves", "myself", "been", "being", "also", "very", "much", "many", "such",
    "than", "ever", "here", "there", "once", "often", "again", "every", "across",
    "of", "in", "on", "at", "by", "is", "are", "be", "am", "was", "were", "it",
    "its", "as", "to", "or", "an", "a", "so", "if", "not", "no", "do", "does",
    "did", "each", "per", "via", "both", "same", "own", "due", "per", "via",
}

_KO_STOPWORDS: Set[str] = {
    "그리고", "그러나", "하지만", "그러면서", "그러므로", "또한", "그러니까", "따라서", "그리고나서",
    "그러면", "그리고도", "그러곤", "그러했지만", "그러할", "그러하다", "그러한", "그런", "이는",
    "이는", "이를", "있는", "있으며", "있습니다", "합니다", "하였다", "하는", "하게", "하고",
    "하여", "하여금", "해서", "하지만", "혹은", "또는", "부터", "까지", "위해", "대한", "다른",
    "모든", "각각", "관련", "경우", "때문", "때문에", "여러", "어떤", "일부", "특히", "다만",
    "즉", "따위", "예를", "예를들어", "수", "등", "및", "것", "그리고", "또", "또한",
    "우리", "너희", "그것", "이것", "저것", "그", "이", "저", "에게", "에서", "으로",
    "로", "에는", "에는", "였다", "이며", "면서", "이라", "이라서",
}

_DOMAIN_STOPWORDS: Set[str] = {
    "document",
    "documents",
    "report",
    "reports",
    "file",
    "files",
    "data",
    "자료",
    "파일",
    "문서",
    "보고서",
    "첨부",
    "자료들",
    "내용",
    "프로젝트",
    "관련자료",
}

_STOPWORDS: Set[str] = {
    word.lower() for word in (*_EN_STOPWORDS, *_KO_STOPWORDS, *_DOMAIN_STOPWORDS)
}



def _split_tokens(source: str) -> List[str]:
    if not source:
        return []
    return [tok for tok in _META_SPLIT_RE.split(source) if tok]


def _remove_stopwords(text: str) -> str:
    if not text:
        return ""
    kept: List[str] = []
    for match in _TOKEN_REGEX.finditer(text):
        token = match.group(0)
        token_norm = token.lower()
        if token_norm in _STOPWORDS:
            continue
        if token_norm.isdigit():
            continue
        if len(set(token_norm)) == 1 and len(token_norm) <= 3:
            continue
        kept.append(token)
    if not kept:
        return text.strip()
    return " ".join(kept)


def _slice_text_by_ratio(source: str, start_char: int, end_char: int, base_len: int) -> str:
    if not source:
        return ""
    if base_len <= 0:
        return source.strip()
    length = len(source)
    start_ratio = max(0.0, min(1.0, float(start_char) / float(base_len)))
    end_ratio = max(start_ratio, min(1.0, float(end_char) / float(base_len)))
    start_idx = int(round(start_ratio * length))
    end_idx = int(round(end_ratio * length))
    if end_idx <= start_idx:
        end_idx = min(length, max(start_idx + 1, end_idx))
    return source[start_idx:end_idx].strip()


def _token_chunk_spans(text: str, *, min_tokens: int, max_tokens: int) -> List[Tuple[int, int, int]]:
    if not text or not text.strip():
        cleaned = (text or "").strip()
        return [(0, len(text), 0)] if cleaned else []

    matches = list(_TOKEN_REGEX.finditer(text))
    total_tokens = len(matches)
    if total_tokens == 0:
        cleaned = text.strip()
        return [(0, len(text), 0)] if cleaned else []
    if total_tokens <= max_tokens:
        return [(0, len(text), total_tokens)]

    spans: List[Tuple[int, int, int]] = []
    start_index = 0
    prev_char = 0
    text_len = len(text)

    while start_index < total_tokens:
        end_index = min(start_index + max_tokens, total_tokens)
        remaining = total_tokens - end_index
        if remaining and remaining < min_tokens:
            end_index = total_tokens
        next_start_char = matches[end_index].start() if end_index < total_tokens else text_len
        span_start = prev_char
        span_end = next_start_char
        token_count = end_index - start_index
        chunk = text[span_start:span_end].strip()
        if chunk:
            spans.append((span_start, span_end, token_count))
        prev_char = next_start_char
        start_index = end_index

    if len(spans) >= 2 and spans[-1][2] < min_tokens:
        prev_start, _prev_end, prev_tokens = spans[-2]
        spans[-2] = (prev_start, spans[-1][1], prev_tokens + spans[-1][2])
        spans.pop()

    if spans and spans[-1][1] < text_len:
        start, _, tokens = spans[-1]
        spans[-1] = (start, text_len, tokens)

    return spans


def _apply_uniform_chunks(
    df: "pd.DataFrame",
    *,
    min_tokens: int = DEFAULT_CHUNK_MIN_TOKENS,
    max_tokens: int = DEFAULT_CHUNK_MAX_TOKENS,
) -> "pd.DataFrame":
    if pd is None or df is None or df.empty or "text" not in df.columns:
        return df

    records = df.to_dict(orient="records")
    chunked: List[Dict[str, Any]] = []

    for record in records:
        base_text = str(record.get("text") or "")
        spans = _token_chunk_spans(base_text, min_tokens=min_tokens, max_tokens=max_tokens)
        if not spans:
            new_rec = dict(record)
            new_rec["chunk_id"] = 1
            new_rec["chunk_count"] = 1
            new_rec["chunk_tokens"] = 0
            preview_source = record.get("text_original") or record.get("text") or ""
            new_rec["text"] = _remove_stopwords(base_text)
            new_rec["text_original"] = preview_source
            new_rec["preview"] = str(preview_source).strip()[:360]
            chunked.append(new_rec)
            continue

        chunk_count = max(1, len(spans))
        base_len = len(base_text)
        original_text = record.get("text_original") or ""

        for idx, (start_char, end_char, token_count) in enumerate(spans, start=1):
            chunk_slice = base_text[start_char:end_char].strip()
            filtered_chunk = _remove_stopwords(chunk_slice)
            if not filtered_chunk:
                filtered_chunk = chunk_slice

            new_rec = dict(record)
            new_rec["chunk_id"] = idx
            new_rec["chunk_count"] = chunk_count
            new_rec["chunk_tokens"] = token_count
            new_rec["text"] = filtered_chunk

            if isinstance(original_text, str) and original_text:
                orig_chunk = _slice_text_by_ratio(original_text, start_char, end_char, base_len)
            else:
                orig_chunk = chunk_slice
            new_rec["text_original"] = orig_chunk
            new_rec["preview"] = (orig_chunk or chunk_slice).strip()[:360]

            chunked.append(new_rec)

    return pd.DataFrame(chunked)


def _time_tokens(epoch: Optional[float]) -> List[str]:
    if not epoch:
        return []
    try:
        dt = datetime.fromtimestamp(float(epoch))
    except Exception:
        return []
    parts = [
        dt.strftime("%Y"),
        dt.strftime("%Y-%m"),
        dt.strftime("%Y-%m-%d"),
        dt.strftime("%B"),
        dt.strftime("%m"),
    ]
    return parts


def _size_bucket(size: Optional[int]) -> Optional[str]:
    if size is None:
        return None
    try:
        size = int(size)
    except (TypeError, ValueError):
        return None
    if size <= 0:
        return None
    if size < 10 * 1024:
        return "size:tiny"
    if size < 1 * 1024 * 1024:
        return "size:small"
    if size < 10 * 1024 * 1024:
        return "size:medium"
    if size < 50 * 1024 * 1024:
        return "size:large"
    return "size:huge"


def _metadata_text(
    path: str,
    ext: str,
    drive: str,
    size: Optional[int] = None,
    mtime: Optional[float] = None,
    ctime: Optional[float] = None,
    owner: Optional[str] = None,
) -> str:
    tokens: List[str] = []
    if path:
        try:
            p = Path(path)
        except Exception:
            p = None
        if p:
            name = p.name
            if name:
                tokens.append(name)
            stem = p.stem
            if stem and stem != name:
                tokens.append(stem)
            tokens.extend(_split_tokens(stem))
            parent_name = p.parent.name if p.parent else ""
            if parent_name:
                tokens.append(parent_name)
                tokens.extend(_split_tokens(parent_name))
        else:
            tokens.append(str(path))
    if ext:
        ext_clean = str(ext).strip()
        if ext_clean:
            tokens.append(ext_clean)
            ext_no_dot = ext_clean.lstrip(".")
            if ext_no_dot:
                tokens.append(ext_no_dot)
    if drive:
        drive_str = str(drive)
        tokens.append(drive_str)
        tokens.extend(_split_tokens(drive_str))
    for epoch in (mtime, ctime):
        tokens.extend(_time_tokens(epoch))
    bucket = _size_bucket(size)
    if bucket:
        tokens.append(bucket)
    if owner:
        tokens.append(str(owner))
        tokens.extend(_split_tokens(str(owner)))

    seen = set()
    normalized: List[str] = []
    for token in tokens:
        cleaned = TextCleaner.clean(str(token)).lower()
        if not cleaned:
            continue
        if cleaned not in seen:
            seen.add(cleaned)
            normalized.append(cleaned)
    return " ".join(normalized)


def _compose_model_text(base_text: str, metadata: str) -> str:
    base_text = base_text or ""
    metadata = metadata or ""
    if metadata and base_text:
        return f"{base_text}\n\n{metadata}"
    if metadata:
        return metadata
    return base_text


def _prepare_text_frame(df: "pd.DataFrame") -> "pd.DataFrame":
    if pd is None or df is None:
        return df
    if df.empty:
        if MODEL_TEXT_COLUMN not in df.columns:
            df[MODEL_TEXT_COLUMN] = pd.Series(dtype=str)
        return df

    for column in ("text", "text_original"):
        if column in df.columns:
            df[column] = df[column].fillna("").astype(str)

    if "text" not in df.columns:
        df["text"] = ""

    paths = df.get("path")
    if paths is None:
        paths = pd.Series([""] * len(df))
    else:
        paths = paths.fillna("").astype(str)

    exts = df.get("ext")
    if exts is None:
        exts = pd.Series([""] * len(df))
    else:
        exts = exts.fillna("").astype(str)

    drives = df.get("drive")
    if drives is None:
        drives = pd.Series([""] * len(df))
    else:
        drives = drives.fillna("").astype(str)

    sizes = df.get("size")
    if sizes is None:
        sizes = pd.Series([0] * len(df))
    else:
        sizes = sizes.fillna(0)

    mtimes = df.get("mtime")
    if mtimes is None:
        mtimes = pd.Series([0.0] * len(df))
    else:
        mtimes = mtimes.fillna(0.0)

    ctimes = df.get("ctime")
    if ctimes is None:
        ctimes = pd.Series([0.0] * len(df))
    else:
        ctimes = ctimes.fillna(0.0)

    owners = df.get("owner")
    if owners is None:
        owners = pd.Series([""] * len(df))
    else:
        owners = owners.fillna("").astype(str)

    base_texts = df["text"].tolist()
    metadata_list = [
        _metadata_text(
            paths.iat[idx],
            exts.iat[idx],
            drives.iat[idx],
            size=sizes.iat[idx],
            mtime=mtimes.iat[idx],
            ctime=ctimes.iat[idx],
            owner=owners.iat[idx],
        )
        for idx in range(len(df))
    ]
    df[MODEL_TEXT_COLUMN] = [
        _compose_model_text(base_texts[idx], metadata_list[idx])
        for idx in range(len(df))
    ]
    return df


def _resolve_kmeans_n_init() -> Union[str, int]:
    """Return MiniBatchKMeans n_init compatible with installed scikit-learn."""
    try:
        parts = (sklearn_version or "0").split(".")
        major = int(parts[0])
        minor = int(parts[1]) if len(parts) > 1 else 0
        if (major, minor) >= (1, 4):
            return "auto"
    except Exception:
        pass
    return 3


# =========================
# Extractors
# =========================
class BaseExtractor:
    exts: Tuple[str,...] = ()
    def can_handle(self, p:Path)->bool: return p.suffix.lower() in self.exts
    def extract(self, p:Path)->Dict[str,Any]: raise NotImplementedError

class HwpExtractor(BaseExtractor):
    exts = (".hwp",)

    def extract(self, p: Path) -> Dict[str, Any]:
        system = platform.system().lower()
        if system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                app = win32com.Dispatch("HWPFrame.HwpObject")
                try:
                    app.Open(str(p))
                    text = app.GetTextFile("TEXT", "") or ""
                    return {
                        "ok": True,
                        "text": TextCleaner.clean(text),
                        "meta": {"engine": "win32com-hwp"},
                    }
                finally:
                    try:
                        app.Quit()
                    except Exception:
                        pass
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"HWP win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass
        if olefile and pyhwp:
            try:
                from pyhwp.hwp5txt import hwp5txt  # type: ignore

                with olefile.OleFileIO(str(p)) as ole:
                    buf = io.StringIO()
                    hwp5txt(ole, buf)
                    text = buf.getvalue()
                cleaned = TextCleaner.clean(text)
                if cleaned:
                    return {
                        "ok": True,
                        "text": cleaned,
                        "meta": {"engine": "pyhwp", "bytes": p.stat().st_size},
                    }
            except Exception as exc:
                return {
                    "ok": False,
                    "text": "",
                    "meta": {"error": f"HWP pyhwp 추출 실패: {exc}"},
                }
        return {
            "ok": False,
            "text": "",
            "meta": {"error": "HWP 추출을 위해서는 Windows + 한/글 환경이 필요합니다."},
        }


class DocDocxExtractor(BaseExtractor):
    exts = (".doc", ".docx")

    def extract(self, p: Path) -> Dict[str, Any]:
        suffix = p.suffix.lower()
        if suffix == ".docx" and docx:
            try:
                document = docx.Document(str(p))
                text = "\n".join(par.text for par in document.paragraphs)
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {"engine": "python-docx", "paras": len(document.paragraphs)},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"DOCX parse failed: {exc}"}}

        system = platform.system().lower()
        if suffix == ".doc" and system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                word = win32com.Dispatch("Word.Application")
                word.Visible = False
                try:
                    doc_obj = word.Documents.Open(str(p), ReadOnly=True)
                    try:
                        text = doc_obj.Content.Text or ""
                    finally:
                        doc_obj.Close(False)
                finally:
                    try:
                        word.Quit()
                    except Exception:
                        pass
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {"engine": "win32com-word"},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"DOC win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass

        return {
            "ok": False,
            "text": "",
            "meta": {"error": "DOC/DOCX 추출을 위해 python-docx 또는 Windows Word가 필요합니다."},
        }

class ExcelLikeExtractor(BaseExtractor):
    exts=(".xlsx",".xls",".xlsm",".xlsb",".xltx",".csv")
    def extract(self, p:Path)->Dict[str,Any]:
        if pd is None:
            return {"ok":False,"text":"","meta":{"error":"pandas required"}}
        try:
            if p.suffix.lower()==".csv":
                df=pd.read_csv(p, nrows=200, encoding="utf-8", engine="python")
                txt=self._df_to_text(df)
                return {"ok":True,"text":txt,"meta":{"engine":"pandas","columns":df.columns.tolist(), "rows_preview":min(200,len(df))}}
            eng = "openpyxl" if p.suffix.lower() in (".xlsx",".xlsm",".xltx") else ("xlrd" if p.suffix.lower()==".xls" else "pyxlsb")
            sheets = pd.read_excel(p, sheet_name=None, nrows=200, engine=eng)
            parts=[]
            for s,df_sheet in sheets.items():
                parts.append(f"[Sheet:{s}]")
                parts.append(" | ".join(map(str, df_sheet.columns.tolist())))
                for _,row in df_sheet.head(50).iterrows():
                    parts.append(" • "+" | ".join(map(lambda x: str(x), row.tolist())))
            return {"ok":True,"text":TextCleaner.clean("\n".join(parts)),"meta":{"engine":"pandas","sheets":list(sheets.keys())}}
        except Exception as e:
            detail = str(e)
            if "openpyxl" in detail.lower():
                detail += " (pip install openpyxl)"
            return {"ok":False,"text":"","meta":{"error":f"excel/csv read failed: {detail}"}}
    @staticmethod
    def _df_to_text(df)->str:
        cols=" | ".join(map(str, df.columns.tolist()))
        rows=[]
        for _,row in df.head(50).iterrows():
            rows.append(" • "+" | ".join(map(lambda x: str(x), row.tolist())))
        return TextCleaner.clean(f"{cols}\n"+"\n".join(rows))

class PdfExtractor(BaseExtractor):
    exts = (".pdf",)

    def extract(self, p: Path) -> Dict[str, Any]:
        if fitz:
            try:
                with fitz.open(str(p)) as doc:
                    page_count = doc.page_count
                    text = "\n".join(page.get_text("text") for page in doc)
                return {
                    "ok": True,
                    "text": TextCleaner.clean(text),
                    "meta": {"engine": "pymupdf", "pages": page_count},
                }
            except Exception:
                pass
        if pdfplumber:
            try:
                with pdfplumber.open(str(p)) as doc:
                    pages = [page.extract_text() or "" for page in doc.pages]
                text = "\n".join(pages)
                cleaned = TextCleaner.clean(text)
                if cleaned:
                    return {
                        "ok": True,
                        "text": cleaned,
                        "meta": {"engine": "pdfplumber", "pages": len(pages)},
                    }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PDF pdfplumber 실패: {exc}"}}
        if pdfminer_extract_text:
            try:
                text = pdfminer_extract_text(str(p))
                return {"ok": True, "text": TextCleaner.clean(text), "meta": {"engine": "pdfminer"}}
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PDF pdfminer 실패: {exc}"}}
        return {"ok": False, "text": "", "meta": {"error": "PDF 추출 엔진이 설치되지 않았습니다."}}


class PptExtractor(BaseExtractor):
    exts = (".ppt", ".pptx")

    def extract(self, p: Path) -> Dict[str, Any]:
        suffix = p.suffix.lower()
        if suffix == ".pptx" and pptx:
            try:
                presentation = pptx.Presentation(str(p))
                texts: List[str] = []
                for idx, slide in enumerate(presentation.slides, 1):
                    parts: List[str] = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text = (shape.text or "").strip()
                            if text:
                                parts.append(text)
                    if parts:
                        texts.append(f"[Slide {idx}] " + " ".join(parts))
                return {
                    "ok": True,
                    "text": TextCleaner.clean("\n".join(texts)),
                    "meta": {"engine": "python-pptx", "slides": len(presentation.slides)},
                }
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PPTX parse failed: {exc}"}}

        system = platform.system().lower()
        if suffix == ".ppt" and system.startswith("win") and win32com:
            com_initialized = False
            try:
                if pythoncom:
                    pythoncom.CoInitialize()
                    com_initialized = True
                powerpoint = win32com.Dispatch("PowerPoint.Application")
                powerpoint.Visible = False
                presentation = powerpoint.Presentations.Open(str(p), WithWindow=False)
                texts: List[str] = []
                try:
                    for slide in presentation.Slides:
                        parts = []
                        for shape in slide.Shapes:
                            has_text = hasattr(shape, "HasTextFrame") and shape.HasTextFrame
                            if has_text and shape.TextFrame.HasText:
                                parts.append(shape.TextFrame.TextRange.Text)
                        if parts:
                            texts.append(" ".join(parts))
                    return {
                        "ok": True,
                        "text": TextCleaner.clean("\n".join(texts)),
                        "meta": {"engine": "win32com-ppt"},
                    }
                finally:
                    presentation.Close()
                    powerpoint.Quit()
            except Exception as exc:
                return {"ok": False, "text": "", "meta": {"error": f"PPT win32com 실패: {exc}"}}
            finally:
                if com_initialized and pythoncom:
                    try:
                        pythoncom.CoUninitialize()
                    except Exception:
                        pass

        return {"ok": False, "text": "", "meta": {"error": "PPT/PPTX 추출을 위해 python-pptx 또는 Windows PowerPoint가 필요합니다."}}


EXTRACTORS = [
    HwpExtractor(),
    DocDocxExtractor(),
    ExcelLikeExtractor(),
    PdfExtractor(),
    PptExtractor(),
]
EXT_MAP={e:ex for ex in EXTRACTORS for e in ex.exts}


# =========================
# 코퍼스 빌더 (번역 기능 수정)
# =========================
@dataclass
class ExtractRecord:
    path: str
    ext: str
    ok: bool
    text: str
    text_original: str
    meta: Dict[str, Any]
    size: Optional[int] = None
    mtime: Optional[float] = None
    ctime: Optional[float] = None
    owner: Optional[str] = None

class CorpusBuilder:
    MAX_TRANSLATE_CHARS = 4000

    def __init__(
        self,
        max_text_chars: int = 200_000,
        progress: bool = True,
        translate: bool = False,
        max_workers: Optional[int] = None,
    ):
        self.max_text_chars = max_text_chars
        self.progress = progress
        self.translate = translate
        self.translator = None
        if translate:
            if GoogleTranslator is None:
                print("⚠️ 경고: 'deep-translator' 라이브러리를 찾을 수 없어 번역 기능이 비활성화됩니다.")
                print("   해결: pip install deep-translator")
            else:
                try:
                    self.translator = GoogleTranslator(source="auto", target="en")
                except Exception as exc:
                    print("⚠️ 경고: 번역기 초기화에 실패해 번역 기능이 비활성화됩니다.")
                    print(f"   상세: {exc}")
        worker_default = max(1, min(8, (os.cpu_count() or 4)))
        self.max_workers = max_workers or worker_default
        if self.translate:
            # 번역 시 외부 API 호출이 순차 처리되도록 워커 1개만 사용
            self.max_workers = 1

    def build(self, file_rows: List[Dict[str, Any]]):
        if pd is None:
            raise RuntimeError("pandas 필요. pip install pandas")

        total = len(file_rows)
        if total == 0:
            print("ℹ️ 신규/변경 문서가 없어 추출을 건너뜁니다.", flush=True)
            return pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

        use_tqdm = self.progress and tqdm is not None
        desc = "📥 Extract & Translate" if self.translate else "📥 Extract"
        bar = tqdm(total=total, desc=desc, unit="file") if use_tqdm else ProgressLine(total, "extracting", update_every=max(1, total // 100 or 1))

        recs: List[Optional[ExtractRecord]] = [None] * total
        with ThreadPoolExecutor(max_workers=max(1, self.max_workers)) as executor:
            future_map = {
                executor.submit(self._extract_one, file_rows[idx]): idx
                for idx in range(total)
            }
            for future in as_completed(future_map):
                idx = future_map[future]
                try:
                    rec = future.result()
                except Exception as exc:
                    row = file_rows[idx]
                    rec = ExtractRecord(
                        path=row.get("path", ""),
                        ext=row.get("ext", ""),
                        ok=False,
                        text="",
                        text_original="",
                        meta={"error": f"extract crash: {exc}"},
                        size=row.get("size"),
                        mtime=row.get("mtime"),
                        ctime=row.get("ctime"),
                        owner=row.get("owner"),
                    )
                recs[idx] = rec
                if use_tqdm:
                    bar.update(1)
                else:
                    bar.update(1)

        if use_tqdm and bar is not None:
            bar.close()
        elif not use_tqdm:
            bar.close()

        records = [r.__dict__ for r in recs if r is not None]
        df = pd.DataFrame(records)
        _prepare_text_frame(df)
        ok = int(df["ok"].sum()) if len(df) > 0 else 0
        fail = int((~df["ok"]).sum()) if len(df) > 0 else 0
        print(f"✅ Extract 완료: ok={ok}, fail={fail}", flush=True)
        return df

    def _extract_one(self, row:Dict[str,Any])->ExtractRecord:
        p=Path(row["path"]); ext=p.suffix.lower()
        ex=EXT_MAP.get(ext)
        if not ex:
            return ExtractRecord(
                str(p),
                ext,
                False,
                "",
                "",
                {"error":"no extractor"},
                row.get("size"),
                row.get("mtime"),
                row.get("ctime"),
                row.get("owner"),
            )
        try:
            out=ex.extract(p)
            original_text=(out.get("text","") or "")[:self.max_text_chars]

            text_for_model = original_text
            if self.translator and original_text.strip():
                text_for_model = self._translate_text(original_text, context=p.name)

            return ExtractRecord(
                str(p),
                ext,
                bool(out.get("ok",False)),
                text_for_model,
                original_text,
                out.get("meta",{}),
                row.get("size"),
                row.get("mtime"),
                row.get("ctime"),
                row.get("owner"),
            )
        except Exception as e:
            return ExtractRecord(
                str(p),
                ext,
                False,
                "",
                "",
                {"error":f"extract crash: {e}"},
                row.get("size"),
                row.get("mtime"),
                row.get("ctime"),
                row.get("owner"),
            )

    def _translate_text(self, text: str, *, context: str) -> str:
        if not self.translator:
            return text
        chunks = self._chunk_text(text, self.MAX_TRANSLATE_CHARS)
        try:
            translated_chunks: List[str] = []
            for chunk in chunks:
                translated = self.translator.translate(chunk)
                translated_chunks.append(self._translated_text(translated, fallback=chunk))
            joined = "\n".join(translated_chunks).strip()
            return joined or text
        except Exception as exc:
            self._log_warning(f"\n[경고] '{context}' 번역 실패. 원본 텍스트 사용. 오류: {exc}")
            return text

    @staticmethod
    def _translated_text(result: Any, *, fallback: str) -> str:
        if isinstance(result, str):
            return result
        text = getattr(result, "text", None)
        if isinstance(text, str) and text.strip():
            return text
        return fallback

    @staticmethod
    def _chunk_text(text: str, limit: int) -> List[str]:
        if len(text) <= limit:
            return [text]
        chunks: List[str] = []
        start = 0
        length = len(text)
        while start < length:
            end = min(length, start + limit)
            split = end
            if end < length:
                for sep in ("\n\n", "\n", " "):
                    idx = text.rfind(sep, start, end)
                    if idx != -1 and idx > start:
                        split = idx + len(sep)
                        break
            if split <= start:
                split = end
            chunks.append(text[start:split])
            start = split
        return chunks

    def _log_warning(self, message: str) -> None:
        if tqdm and self.progress:
            tqdm.write(message)
        else:
            print(message)

    @staticmethod
    def save(df, out_path:Path):
        out_path.parent.mkdir(parents=True, exist_ok=True)
        ext = out_path.suffix.lower()
        if ext == ".parquet":
            engine_kwargs = {}
            engine_label = PARQUET_ENGINE or "auto"
            if PARQUET_ENGINE:
                engine_kwargs["engine"] = PARQUET_ENGINE
            try:
                df.to_parquet(out_path, index=False, **engine_kwargs)
                print(f"✅ Parquet 저장({engine_label}): {out_path}")
                return
            except Exception as e:
                csv_path = out_path.with_suffix(".csv")
                df.to_csv(csv_path, index=False, encoding="utf-8")
                print(
                    f"⚠️ Parquet 엔진 실패({engine_label}) → CSV로 저장: {csv_path}\n"
                    f"   상세: {e}"
                )
                return
        df.to_csv(out_path, index=False, encoding="utf-8")
        print(f"✅ CSV 저장: {out_path}")


def _load_existing_corpus(path: Path) -> Optional["pd.DataFrame"]:
    if pd is None:
        return None
    candidates = [path]
    suffix = path.suffix.lower()
    if suffix == ".parquet":
        candidates.append(path.with_suffix(".csv"))
    elif suffix == ".csv":
        candidates.append(path.with_suffix(".parquet"))

    for candidate in candidates:
        if not candidate.exists():
            continue
        try:
            if candidate.suffix.lower() == ".parquet":
                engine_kwargs = {}
                if PARQUET_ENGINE:
                    engine_kwargs["engine"] = PARQUET_ENGINE
                return pd.read_parquet(candidate, **engine_kwargs)
            return pd.read_csv(candidate)
        except Exception as exc:
            engine_label = PARQUET_ENGINE or "auto"
            print(
                f"⚠️ 기존 코퍼스 로드 실패 ({candidate}, engine={engine_label}): {exc}",
                flush=True,
            )
    return None


def _is_cache_fresh(cached: Dict[str, Any], row: Dict[str, Any]) -> bool:
    if not cached.get("ok"):
        return False
    if not cached.get("text"):
        return False
    try:
        cached_size = int(cached.get("size", -1))
        row_size = int(row.get("size", -1))
    except (TypeError, ValueError):
        return False
    if cached_size != row_size:
        return False
    try:
        cached_mtime = float(cached.get("mtime", 0.0))
        row_mtime = float(row.get("mtime", 0.0))
    except (TypeError, ValueError):
        return False
    if abs(cached_mtime - row_mtime) > 1.0:
        return False
    return True


def _split_cache(
    file_rows: List[Dict[str, Any]],
    existing_df: Optional["pd.DataFrame"],
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    if existing_df is None or "path" not in existing_df.columns:
        return list(file_rows), []

    cache_map: Dict[str, Dict[str, Any]] = {}
    for rec in existing_df.to_dict(orient="records"):
        cache_map[rec.get("path", "")] = rec

    to_process: List[Dict[str, Any]] = []
    reused: List[Dict[str, Any]] = []
    for row in file_rows:
        path = row.get("path")
        cached = cache_map.get(path)
        if cached and _is_cache_fresh(cached, row):
            cached_copy = dict(cached)
            cached_copy["size"] = row.get("size", cached_copy.get("size"))
            cached_copy["mtime"] = row.get("mtime", cached_copy.get("mtime"))
            reused.append(cached_copy)
        else:
            to_process.append(row)
    return to_process, reused


# =========================
# 토픽 모델
# =========================
@dataclass
class TrainConfig:
    max_features: int = 50_000
    n_components: int = DEFAULT_N_COMPONENTS
    n_clusters: int = 30
    ngram_range: Tuple[int, int] = (1, 2)
    min_df: int = 2
    max_df: float = 0.8
    use_sentence_transformer: bool = True
    embedding_model: str = DEFAULT_EMBED_MODEL
    embedding_batch_size: int = 32

class TopicModel:
    def __init__(self, cfg:TrainConfig):
        if any(x is None for x in (TfidfVectorizer, TruncatedSVD, MiniBatchKMeans, Pipeline)):
            raise RuntimeError("scikit-learn 필요. pip install scikit-learn joblib")
        self.cfg=cfg
        self.pipeline:Optional[Pipeline]=None
        self._kmeans_n_init = _resolve_kmeans_n_init()

    def fit(self, df, text_col="text"):
        texts=(df[text_col].fillna("").astype(str)).tolist()
        print("🧠 학습 준비: TF-IDF → SVD → KMeans", flush=True)
        spin=Spinner(prefix="  학습 중")
        spin.start()
        try:
            self.pipeline = Pipeline(steps=[
                ("tfidf", TfidfVectorizer(
                    token_pattern=TOKEN_PATTERN,
                    ngram_range=self.cfg.ngram_range,
                    max_features=self.cfg.max_features,
                    min_df=self.cfg.min_df,
                    max_df=self.cfg.max_df,
                )),
                ("svd", TruncatedSVD(n_components=self.cfg.n_components, random_state=42)),
                ("kmeans", MiniBatchKMeans(n_clusters=self.cfg.n_clusters, random_state=42, batch_size=2048, n_init=self._kmeans_n_init)),
            ])
            t0=time.time()
            self.pipeline.fit(texts)
            t1=time.time()
        finally:
            spin.stop()
        print(f"✅ 학습 완료 (docs={len(texts):,}, {t1-t0:.1f}s)", flush=True)
        return self

    def predict(self, df, text_col="text")->List[int]:
        texts=(df[text_col].fillna("").astype(str)).tolist()
        return self.pipeline.predict(texts)

    def transform(self, df, text_col="text"):
        texts=(df[text_col].fillna("").astype(str)).tolist()
        X=self.pipeline.named_steps["tfidf"].transform(texts)
        Z=self.pipeline.named_steps["svd"].transform(X)
        return Z

    def save(self, path:Path):
        if joblib is None: raise RuntimeError("joblib 필요")
        path.parent.mkdir(parents=True, exist_ok=True)
        joblib.dump({"cfg":self.cfg,"pipeline":self.pipeline}, path)


class SentenceBertModel:
    def __init__(self, cfg: TrainConfig):
        if SentenceTransformer is None:
            raise RuntimeError(
                "sentence-transformers 라이브러리가 필요합니다. pip install sentence-transformers"
            )
        self.cfg = cfg
        self.model_name = cfg.embedding_model or DEFAULT_EMBED_MODEL
        print(f"🧠 Sentence-BERT 준비: {self.model_name}", flush=True)
        self._encoder = SentenceTransformer(self.model_name)
        self.embedding_dim = int(self._encoder.get_sentence_embedding_dimension())
        self.cluster_model: Optional[MiniBatchKMeans] = None
        self.cluster_labels_: Optional[np.ndarray] = None
        self._kmeans_n_init = _resolve_kmeans_n_init()

    def encode(self, texts: List[str], *, show_progress: bool = False) -> np.ndarray:
        if not texts:
            return np.zeros((0, self.embedding_dim), dtype=np.float32)
        embeddings = self._encoder.encode(
            texts,
            batch_size=max(1, int(self.cfg.embedding_batch_size)),
            show_progress_bar=show_progress,
            convert_to_numpy=True,
            normalize_embeddings=False,
        )
        if isinstance(embeddings, list):
            embeddings = np.asarray(embeddings, dtype=np.float32)
        return np.asarray(embeddings, dtype=np.float32)

    def fit(self, df, text_col: str = "text") -> np.ndarray:
        texts = (df[text_col].fillna("").astype(str)).tolist()
        show_progress = tqdm is not None and len(texts) > 1000
        embeddings = self.encode(texts, show_progress=show_progress)

        can_cluster = (
            MiniBatchKMeans is not None
            and self.cfg.n_clusters > 0
            and embeddings.shape[0] >= max(10, self.cfg.n_clusters)
        )
        if can_cluster:
            print("🔖 클러스터링: MiniBatchKMeans", flush=True)
            self.cluster_model = MiniBatchKMeans(
                n_clusters=self.cfg.n_clusters,
                random_state=42,
                batch_size=2048,
                n_init=self._kmeans_n_init,
            )
            self.cluster_model.fit(embeddings)
            try:
                labels = self.cluster_model.labels_
            except AttributeError:
                labels = self.cluster_model.predict(embeddings)
            self.cluster_labels_ = np.asarray(labels, dtype=np.int32)
        else:
            self.cluster_model = None
            self.cluster_labels_ = None
            if MiniBatchKMeans is None:
                print("⚠️ scikit-learn MiniBatchKMeans 미설치로 토픽 라벨링을 건너뜁니다.", flush=True)
            elif embeddings.shape[0] < max(10, self.cfg.n_clusters):
                print("ℹ️ 문서 수가 적어 토픽 클러스터링을 건너뜁니다.", flush=True)
        return embeddings

    def predict(self, embeddings: np.ndarray) -> np.ndarray:
        if self.cluster_model is None:
            raise RuntimeError("클러스터링 모델이 초기화되지 않았습니다.")
        labels = self.cluster_model.predict(embeddings)
        return np.asarray(labels, dtype=np.int32)

    def save(self, path: Path) -> None:
        if joblib is None:
            raise RuntimeError("joblib 필요. pip install joblib")
        path.parent.mkdir(parents=True, exist_ok=True)
        payload: Dict[str, Any] = {
            "version": 2,
            "model_type": MODEL_TYPE_SENTENCE_TRANSFORMER,
            "model_name": self.model_name,
            "embedding_dim": self.embedding_dim,
            "train_config": self.cfg,
        }
        if self.cluster_model is not None:
            payload["cluster_model"] = self.cluster_model
        joblib.dump(payload, path)


# =========================
# 파이프라인 실행 (메인 함수)
# =========================
def run_step2(file_rows:List[Dict[str,Any]],
              out_corpus:Path=Path("./corpus.parquet"),
              out_model:Path=Path("./topic_model.joblib"),
              cfg:TrainConfig=TrainConfig(),
              use_tqdm:bool=True,
              translate:bool=False):
    global tqdm
    original_tqdm = tqdm
    if not use_tqdm:
        tqdm=None

    try:
        print("=== Step 2 시작: 내용 추출 & 학습 === (번역: " + ("활성" if translate else "비활성") + ")", flush=True)
        t_all=time.time()
        if pd is None:
            raise RuntimeError("pandas 필요")

        existing_df = _load_existing_corpus(out_corpus)
        to_process, reused_records = _split_cache(file_rows, existing_df)

        reused_count = len(reused_records)
        process_count = len(to_process)
        total_count = len(file_rows)
        if total_count:
            print(
                f"🗃️ 캐시 재사용: {reused_count:,} | 신규/변경 추출: {process_count:,} | 총합: {total_count:,}",
                flush=True,
            )
        else:
            print("📂 처리할 파일이 없습니다.", flush=True)

        if process_count == 0 and reused_count == total_count and total_count > 0:
            df = pd.DataFrame(reused_records)
            df = _apply_uniform_chunks(
                df,
                min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
                max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
            )
            _prepare_text_frame(df)
            order_map = {row["path"]: idx for idx, row in enumerate(file_rows)}
            if "path" in df.columns:
                df["_order"] = df["path"].map(order_map)
                df = df.sort_values("_order").drop(columns=["_order"]).reset_index(drop=True)
            CorpusBuilder.save(df, out_corpus)
            if out_model.exists():
                try:
                    os.utime(out_model, None)
                except OSError:
                    pass
            print("✨ 변경된 문서가 없어 기존 모델을 유지합니다.", flush=True)
            return df, None

        cb = CorpusBuilder(max_text_chars=200_000, progress=use_tqdm, translate=translate)
        df_new = cb.build(to_process) if process_count else pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

        frames: List["pd.DataFrame"] = []
        if reused_count:
            frames.append(pd.DataFrame(reused_records))
        if not df_new.empty:
            frames.append(df_new)
        if frames:
            df = pd.concat(frames, ignore_index=True)
        else:
            df = pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

        if pd is None:
            raise RuntimeError("pandas 필요")

        df = _apply_uniform_chunks(
            df,
            min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
            max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
        )
        _prepare_text_frame(df)

        order_map = {row["path"]: idx for idx, row in enumerate(file_rows)} if file_rows else {}
        if "path" in df.columns and order_map:
            df["_order"] = df["path"].map(order_map)
            df = df.sort_values("_order").drop(columns=["_order"]).reset_index(drop=True)

        if "ok" in df.columns:
            df["ok"] = df["ok"].apply(lambda v: bool(v) if isinstance(v, bool) else str(v).strip().lower() in {"true", "1", "yes"})
        if "topic" in df.columns:
            df = df.drop(columns=["topic"])

        text_col = MODEL_TEXT_COLUMN if MODEL_TEXT_COLUMN in df.columns else "text"
        text_mask = df[text_col].fillna("").str.len() > 0
        train_df = df[df["ok"] & text_mask].copy()
        if not train_df.empty:
            _prepare_text_frame(train_df)
        print(f"🧹 학습 대상 문서: {len(train_df):,}/{len(df):,}", flush=True)
        if len(train_df) == 0:
            CorpusBuilder.save(df, out_corpus)
            print(f"⚠️ 유효 텍스트 없음. 코퍼스만 저장: {out_corpus}", flush=True)
            return df, None

        topics_df = None
        model_obj: Optional[Any] = None

        if cfg.use_sentence_transformer and SentenceTransformer is not None:
            try:
                semantic_model = SentenceBertModel(cfg)
                embeddings = semantic_model.fit(train_df, text_col=text_col)
                print(
                    f"✅ Sentence-BERT 임베딩 완료 (docs={embeddings.shape[0]:,}, dim={semantic_model.embedding_dim})",
                    flush=True,
                )
                if semantic_model.cluster_labels_ is not None:
                    train_df["topic"] = semantic_model.cluster_labels_
                    topics_df = train_df[["path", "topic"]].copy()
                model_obj = semantic_model
            except Exception as exc:
                print(f"⚠️ Sentence-BERT 학습 실패로 TF-IDF 백업 경로를 사용합니다: {exc}", flush=True)
        elif cfg.use_sentence_transformer and SentenceTransformer is None:
            print("⚠️ sentence-transformers 미설치로 TF-IDF 백업 경로를 사용합니다.", flush=True)

        if model_obj is None:
            tm = TopicModel(cfg)
            tm.fit(train_df, text_col=text_col)
            labels = tm.predict(train_df, text_col=text_col)
            train_df["topic"] = labels
            topics_df = train_df[["path", "topic"]].copy()
            model_obj = tm

        if topics_df is not None:
            df = df.merge(topics_df, on="path", how="left")

        CorpusBuilder.save(df, out_corpus)

        if isinstance(model_obj, SentenceBertModel):
            model_obj.save(out_model)
        elif isinstance(model_obj, TopicModel) and joblib:
            model_obj.save(out_model)

        dt_all = time.time() - t_all
        print(f"💾 저장 완료: corpus → {out_corpus} | model → {out_model}", flush=True)
        print(f"🎉 Step 2 종료 (총 {dt_all:.1f}s)", flush=True)
        return df, model_obj
    finally:
        tqdm = original_tqdm


def update_corpus_file(
    new_records: "pd.DataFrame",
    corpus_path: Path,
) -> "pd.DataFrame":
    """Merge `new_records` into the persisted corpus and return the updated frame."""
    if pd is None:
        raise RuntimeError("pandas 필요. pip install pandas")

    existing = _load_existing_corpus(corpus_path)
    if existing is None or existing.empty:
        combined = new_records.copy()
    else:
        if "path" in existing.columns and "path" in new_records.columns:
            paths_to_replace = set(new_records["path"].astype(str).tolist())
            mask = ~existing["path"].astype(str).isin(paths_to_replace)
            combined = pd.concat([existing[mask], new_records], ignore_index=True)
        else:
            combined = pd.concat([existing, new_records], ignore_index=True)

    combined = _apply_uniform_chunks(
        combined,
        min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
        max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
    )
    _prepare_text_frame(combined)
    CorpusBuilder.save(combined, corpus_path)
    return combined


def remove_from_corpus(paths: List[str], corpus_path: Path) -> "pd.DataFrame":
    """Remove documents whose paths match `paths` from the persisted corpus."""
    if pd is None:
        raise RuntimeError("pandas 필요. pip install pandas")

    existing = _load_existing_corpus(corpus_path)
    if existing is None or existing.empty:
        return pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

    to_drop = {str(p) for p in paths}
    if "path" not in existing.columns:
        return existing

    filtered = existing[~existing["path"].astype(str).isin(to_drop)].copy()
    filtered = _apply_uniform_chunks(
        filtered,
        min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
        max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
    )
    _prepare_text_frame(filtered)
    CorpusBuilder.save(filtered, corpus_path)
    return filtered
