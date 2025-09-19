from pathlib import Path
from typing import List
import sys

import pytest

sys.path.insert(0, str(Path(__file__).resolve().parents[1]))

try:  # pragma: no cover - optional dependency guard
    import pandas as pd
except ImportError:  # pragma: no cover - handled via skip
    pd = None

try:  # pragma: no cover
    from docx import Document
except ImportError:  # pragma: no cover
    Document = None

try:  # pragma: no cover
    from pptx import Presentation
except ImportError:  # pragma: no cover
    Presentation = None

try:  # pragma: no cover
    from reportlab.pdfgen import canvas
except ImportError:  # pragma: no cover
    canvas = None

if pd is None or Document is None or Presentation is None or canvas is None:
    pytest.skip("smoke test requires pandas/docx/pptx/reportlab", allow_module_level=True)

from pipeline import run_step2, TrainConfig
from retriever import Retriever


def _write_pdf(path: Path, text: str) -> None:
    c = canvas.Canvas(str(path))
    c.drawString(72, 720, text)
    c.save()


def _write_docx(path: Path, text: str) -> None:
    doc = Document()
    doc.add_paragraph(text)
    doc.save(str(path))


def _write_pptx(path: Path, title: str, body: str) -> None:
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = body
    prs.save(str(path))


def _write_xlsx(path: Path) -> None:
    df = pd.DataFrame(
        {
            "Task": ["Kickoff", "Review"],
            "Owner": ["Alice", "Bob"],
            "Status": ["Open", "Done"],
        }
    )
    df.to_excel(path, index=False)


def _write_csv(path: Path) -> None:
    df = pd.DataFrame(
        {
            "Name": ["Alice", "Bob"],
            "Email": ["alice@example.com", "bob@example.com"],
        }
    )
    df.to_csv(path, index=False, encoding="utf-8-sig")


def _write_txt(path: Path, text: str) -> None:
    path.write_text(text, encoding="utf-8-sig")


def _make_sample_files(base: Path) -> List[Path]:
    files = []
    pdf_path = base / "reports" / "budget.pdf"
    _write_pdf(pdf_path, "Budget report and forecast")
    files.append(pdf_path)

    docx_path = base / "memos" / "meeting.docx"
    _write_docx(docx_path, "Meeting schedule summary for the team")
    files.append(docx_path)

    pptx_path = base / "slides" / "overview.pptx"
    _write_pptx(pptx_path, "Project overview", "Timeline and milestones")
    files.append(pptx_path)

    xlsx_path = base / "tables" / "tasks.xlsx"
    _write_xlsx(xlsx_path)
    files.append(xlsx_path)

    csv_path = base / "tables" / "contacts.csv"
    _write_csv(csv_path)
    files.append(csv_path)

    txt_path = base / "notes" / "ideas.txt"
    _write_txt(txt_path, "Brainstorming ideas for next sprint")
    files.append(txt_path)

    return files


def _make_file_rows(paths: List[Path]) -> List[dict]:
    rows = []
    for p in paths:
        stat = p.stat()
        rows.append(
            {
                "path": str(p),
                "size": stat.st_size,
                "mtime": stat.st_mtime,
                "ext": p.suffix.lower(),
                "drive": p.anchor,
            }
        )
    return rows


def test_run_step2_extracts_and_retrieves(tmp_path):
    files = _make_sample_files(tmp_path)
    rows = _make_file_rows(files)

    corpus_path = tmp_path / "data" / "corpus.parquet"
    model_path = tmp_path / "data" / "topic_model.joblib"
    cache_dir = tmp_path / "cache"

    cfg = TrainConfig(
        max_features=512,
        n_components=8,
        n_clusters=4,
        ngram_range=(1, 1),
        min_df=1,
        max_df=1.0,
    )

    df, _ = run_step2(rows, out_corpus=corpus_path, out_model=model_path, cfg=cfg, use_tqdm=False, translate=False)

    assert len(df) == len(rows)
    ok = df["ok"].sum()
    assert ok / len(rows) >= 0.6

    retr = Retriever(model_path=model_path, corpus_path=corpus_path, cache_dir=cache_dir)
    retr.ready(rebuild=True)
    hits = retr.search("project overview timeline", top_k=3)
    assert hits
    assert any(Path(hit["path"]).suffix.lower() == ".pptx" for hit in hits)
