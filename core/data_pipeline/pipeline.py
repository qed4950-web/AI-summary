# pipeline.py  (Step2: 추출 + 학습)
import importlib
import subprocess
import json
import tempfile
import io
import math
import os
import platform
import re
import sys
import threading
import time
import hashlib
from datetime import datetime
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path
from dataclasses import dataclass, replace
from typing import Optional, Dict, Any, List, Tuple, Union, Set

from core.config.paths import MODELS_DIR
from core.data_pipeline.custom_metadata import get_metadata_for_path
from core.data_pipeline.cache_manager import ChunkCache, SQLiteChunkCache
from core.data_pipeline.incremental import (
    load_scan_state,
    filter_incremental_rows,
    update_scan_state,
    save_scan_state,
)
from core.data_pipeline.embedder import AsyncSentenceEmbedder
from core.data_pipeline.chunking_v2 import SemanticChunker

# Backward compatibility: extractors extracted to extractors.py
from core.data_pipeline.extractors import (
    BaseExtractor,
    HwpExtractor,
    DocDocxExtractor,
    ExcelLikeExtractor,
    PdfExtractor,
    PptExtractor,
    PlainTextExtractor,
    CodeExtractor,
    EXTRACTORS,
    EXT_MAP,
)
from core.data_pipeline.evaluate import evaluate_embeddings

# Backward compatibility: models extracted to models.py
from core.data_pipeline.models import (
    TrainConfig,
    TopicModel,
    SentenceBertModel,
    default_train_config,
    _fit_sentence_transformer_chunked,
    _chunked_sentence_embeddings,
    _run_embed_chunk_subprocess,
    _resolve_embed_dtype,
    DEFAULT_EMBED_MODEL,
    DEFAULT_N_COMPONENTS,
    MODEL_TYPE_SENTENCE_TRANSFORMER,
)

# Backward compatibility: corpus_builder extracted to corpus_builder.py
from core.data_pipeline.corpus_builder import (
    ExtractRecord,
    CorpusBuilder,
    _load_existing_corpus,
    _is_cache_fresh,
    _split_cache,
    _collect_existing_rows,
)

import numpy as np

# New Utils
from core.utils.cli_ui import Spinner, ProgressLine
from core.utils.nlp import TextCleaner, split_tokens as _split_tokens_util
from core.utils.stopwords import STOPWORDS as _STOPWORDS


# ---- 선택 의존성(있으면 사용) ----
try:
    import pandas as pd
except Exception:
    pd = None
PARQUET_ENGINE: Optional[str] = None
if pd is not None:
    for candidate in ("fastparquet", "pyarrow"):
        try:
            importlib.import_module(candidate)
            PARQUET_ENGINE = candidate
            break
        except ImportError:
            continue
try:
    from deep_translator import GoogleTranslator
except Exception:
    GoogleTranslator = None
try:
    import docx
except Exception:
    docx = None
try:
    import pptx
except Exception:
    pptx = None
try:
    from pdfminer.high_level import extract_text as pdfminer_extract_text
except Exception:
    pdfminer_extract_text = None
try:
    import win32com.client
except Exception:
    win32com = None
try:
    import pythoncom
except Exception:
    pythoncom = None
try:
    import fitz  # PyMuPDF
except Exception:
    fitz = None
try:
    import joblib
except Exception:
    joblib = None
try:
    import pdfplumber
except Exception:
    pdfplumber = None
try:
    from sklearn.feature_extraction.text import TfidfVectorizer
    from sklearn.decomposition import TruncatedSVD
    from sklearn.cluster import MiniBatchKMeans
    from sklearn.pipeline import Pipeline
    from sklearn.preprocessing import FunctionTransformer
    from sklearn import __version__ as sklearn_version
except Exception:
    TfidfVectorizer = TruncatedSVD = MiniBatchKMeans = Pipeline = FunctionTransformer = None
    sklearn_version = "0"

try:
    from tqdm import tqdm
except Exception:
    tqdm = None

try:
    from sentence_transformers import SentenceTransformer
except Exception:
    SentenceTransformer = None
try:
    import olefile
except Exception:
    olefile = None
try:
    import pyhwp
except Exception:
    pyhwp = None


# =========================
# 콘솔 진행도 유틸
# =========================
# [Refactor] UI classes and TextCleaner moved to core.utils


TOKEN_PATTERN = r'(?u)(?:[가-힣]{1,}|[A-Za-z0-9]{2,})'

# 고정된 SVD 차원 수. Index/모델 불일치를 막기 위해 한곳에서 정의한다.
DEFAULT_N_COMPONENTS = 128
MODEL_TEXT_COLUMN = "text_model"
_META_SPLIT_RE = re.compile(r"[^0-9A-Za-z가-힣]+")
def _default_embed_model() -> str:
    env_model = os.getenv("DEFAULT_EMBED_MODEL")
    if env_model:
        return env_model
    if platform.system() == "Darwin":
        # Prefer the bundled multilingual-e5-small copy on macOS for stability
        return "models--intfloat--multilingual-e5-small"
    return "BAAI/bge-m3"


DEFAULT_EMBED_MODEL = _default_embed_model()
MODEL_TYPE_SENTENCE_TRANSFORMER = "sentence-transformer"

def _normalize_hf_model_id(value: str) -> str:
    """Convert cache-style ids like `models--org--repo` into `org/repo`."""
    raw = (value or "").strip()
    if not raw:
        return raw
    if raw.startswith("models--") and "/" not in raw:
        parts = raw.split("--")
        if len(parts) >= 3:
            org = parts[1].strip()
            repo = "--".join(parts[2:]).strip()
            if org and repo:
                return f"{org}/{repo}"
    return raw


def _resolve_sentence_transformer_location(model_name: str) -> str:
    """Prefer local model snapshots under `models/sentence_transformers/` when available."""
    base_dir = MODELS_DIR / "sentence_transformers"
    if not base_dir.exists():
        return model_name

    direct = base_dir / model_name
    if direct.exists():
        snapshots = direct / "snapshots"
        if snapshots.exists():
            candidates = sorted(
                snapshots.iterdir(),
                key=lambda item: item.stat().st_mtime,
                reverse=True,
            )
            for candidate in candidates:
                if any(
                    (candidate / marker).exists()
                    for marker in ("config.json", "modules.json", "config_sentence_transformers.json")
                ):
                    return str(candidate)
        if any(
            (direct / marker).exists()
            for marker in ("config.json", "modules.json", "config_sentence_transformers.json")
        ):
            return str(direct)

    cache_root = base_dir / f"models--{model_name.replace('/', '--')}"
    if cache_root.exists():
        snapshots = cache_root / "snapshots"
        if snapshots.exists():
            candidates = sorted(
                snapshots.iterdir(),
                key=lambda item: item.stat().st_mtime,
                reverse=True,
            )
            for candidate in candidates:
                if any(
                    (candidate / marker).exists()
                    for marker in ("config.json", "modules.json", "config_sentence_transformers.json")
                ):
                    return str(candidate)
        if any(
            (cache_root / marker).exists()
            for marker in ("config.json", "modules.json", "config_sentence_transformers.json")
        ):
            return str(cache_root)

    return model_name

DEFAULT_CHUNK_MIN_TOKENS = 200
DEFAULT_CHUNK_MAX_TOKENS = 500

EMBED_DTYPE_ENV = "INFOPILOT_EMBED_DTYPE"
_VALID_EMBED_DTYPES = {"auto", "fp16", "fp32"}
CACHE_BACKEND_ENV = "INFOPILOT_CACHE_BACKEND"
_VALID_CACHE_BACKENDS = {"json", "sqlite"}


def _sanitize_embed_dtype(value: Optional[str]) -> Optional[str]:
    if value is None:
        return None
    normalized = str(value).strip().lower()
    return normalized if normalized in _VALID_EMBED_DTYPES else None


def _sanitize_cache_backend(value: Optional[str]) -> str:
    if not value:
        return "json"
    normalized = str(value).strip().lower()
    return normalized if normalized in _VALID_CACHE_BACKENDS else "json"


def _create_chunk_cache(path: Path) -> ChunkCache:
    backend = _sanitize_cache_backend(os.getenv(CACHE_BACKEND_ENV))
    actual_path = path
    if backend == "sqlite":
        if actual_path.suffix.lower() == ".json":
            actual_path = actual_path.with_suffix(".sqlite")
        elif not actual_path.name.endswith(".sqlite"):
            actual_path = actual_path.with_name(actual_path.name + ".sqlite")
    if backend == "sqlite":
        print(f"⚙️ Chunk cache: SQLite backend → {actual_path}", flush=True)
        return SQLiteChunkCache(actual_path)
    if backend != "json":
        print(f"⚠️ 지원하지 않는 캐시 백엔드 '{backend}' → json으로 대체합니다.", flush=True)
    return ChunkCache(path)

_TOKEN_REGEX = re.compile(TOKEN_PATTERN)


def _hash_text(text: str) -> str:
    if not text:
        return ""
    return hashlib.sha1(text.encode("utf-8", "ignore")).hexdigest()

# [Refactor] Stopwords moved to core.utils




def _split_tokens(source: str) -> List[str]:
    if not source:
        return []
    return _split_tokens_util(source)


def _remove_stopwords(text: str) -> str:
    if not text:
        return ""
    kept: List[str] = []
    for match in _TOKEN_REGEX.finditer(text):
        token = match.group(0)
        token_norm = token.lower()
        if token_norm in _STOPWORDS:
            continue
        if token_norm.isdigit():
            continue
        if len(set(token_norm)) == 1 and len(token_norm) <= 3:
            continue
        kept.append(token)
    if not kept:
        return text.strip()
    return " ".join(kept)


def _slice_text_by_ratio(source: str, start_char: int, end_char: int, base_len: int) -> str:
    if not source:
        return ""
    if base_len <= 0:
        return source.strip()
    length = len(source)
    start_ratio = max(0.0, min(1.0, float(start_char) / float(base_len)))
    end_ratio = max(start_ratio, min(1.0, float(end_char) / float(base_len)))
    start_idx = int(round(start_ratio * length))
    end_idx = int(round(end_ratio * length))
    if end_idx <= start_idx:
        end_idx = min(length, max(start_idx + 1, end_idx))
    return source[start_idx:end_idx].strip()


def _token_chunk_spans(text: str, *, min_tokens: int, max_tokens: int) -> List[Tuple[int, int, int]]:
    if not text or not text.strip():
        cleaned = (text or "").strip()
        return [(0, len(text), 0)] if cleaned else []

    matches = list(_TOKEN_REGEX.finditer(text))
    total_tokens = len(matches)
    if total_tokens == 0:
        cleaned = text.strip()
        return [(0, len(text), 0)] if cleaned else []
    if total_tokens <= max_tokens:
        return [(0, len(text), total_tokens)]

    spans: List[Tuple[int, int, int]] = []
    start_index = 0
    prev_char = 0
    text_len = len(text)

    while start_index < total_tokens:
        end_index = min(start_index + max_tokens, total_tokens)
        remaining = total_tokens - end_index
        if remaining and remaining < min_tokens:
            end_index = total_tokens
        next_start_char = matches[end_index].start() if end_index < total_tokens else text_len
        span_start = prev_char
        span_end = next_start_char
        token_count = end_index - start_index
        chunk = text[span_start:span_end].strip()
        if chunk:
            spans.append((span_start, span_end, token_count))
        prev_char = next_start_char
        start_index = end_index

    if len(spans) >= 2 and spans[-1][2] < min_tokens:
        prev_start, _prev_end, prev_tokens = spans[-2]
        spans[-2] = (prev_start, spans[-1][1], prev_tokens + spans[-1][2])
        spans.pop()

    if spans and spans[-1][1] < text_len:
        start, _, tokens = spans[-1]
        spans[-1] = (start, text_len, tokens)

    return spans


def _token_chunk_spans_with_overlap(
    text: str,
    *,
    min_tokens: int,
    max_tokens: int,
    overlap_tokens: int,
) -> List[Tuple[int, int, int]]:
    min_tokens = max(1, int(min_tokens))
    max_tokens = max(min_tokens, int(max_tokens))
    overlap = max(0, int(overlap_tokens))

    matches = list(_TOKEN_REGEX.finditer(text))
    if not matches:
        cleaned = (text or "").strip()
        return [(0, len(text), 0)] if cleaned else []

    total_tokens = len(matches)
    if total_tokens <= max_tokens:
        return [(0, len(text), total_tokens)]

    spans: List[Tuple[int, int, int]] = []
    start_index = 0
    text_len = len(text)
    while start_index < total_tokens:
        end_index = min(start_index + max_tokens, total_tokens)
        remaining = total_tokens - end_index
        if remaining and remaining < min_tokens:
            end_index = total_tokens

        start_char = matches[start_index].start()
        end_char = matches[end_index - 1].end() if end_index > start_index else min(text_len, start_char)
        chunk = text[start_char:end_char].strip()
        if chunk:
            spans.append((start_char, end_char, end_index - start_index))

        if end_index >= total_tokens:
            break
        next_start = end_index - overlap if overlap else end_index
        if next_start <= start_index:
            next_start = end_index
        start_index = next_start

    if spans and spans[-1][1] < text_len:
        start, _, tokens = spans[-1]
        spans[-1] = (start, text_len, tokens)

    return spans


_MD_HEADING_RE = re.compile(r"(?m)^(#{1,6})\s+(.+?)\s*$")
_MD_NUMBERED_HEADING_RE = re.compile(r"(?m)^\s*(\d+(?:\.\d+)*[\).])\s+(.+?)\s*$")


def _iter_markdown_sections(text: str) -> List[Tuple[int, int, str]]:
    """Return (start_char, end_char, heading_title) for markdown-ish sections."""
    if not text:
        return []
    headings: List[Tuple[int, str]] = []
    for match in _MD_HEADING_RE.finditer(text):
        title = (match.group(2) or "").strip()
        headings.append((match.start(), title))
    for match in _MD_NUMBERED_HEADING_RE.finditer(text):
        title = (match.group(2) or "").strip()
        headings.append((match.start(), title))
    if not headings:
        return [(0, len(text), "")]

    headings.sort(key=lambda item: item[0])
    sections: List[Tuple[int, int, str]] = []
    for idx, (start, title) in enumerate(headings):
        end = headings[idx + 1][0] if idx + 1 < len(headings) else len(text)
        sections.append((start, end, title))
    if sections and sections[0][0] > 0:
        sections.insert(0, (0, sections[0][0], ""))
    return [(s, e, t) for (s, e, t) in sections if s < e]


def _is_markdown_record(record: Dict[str, Any]) -> bool:
    ext = str(record.get("ext") or "").lower()
    if ext == ".md":
        return True
    meta = record.get("meta")
    if isinstance(meta, dict) and str(meta.get("format") or "").lower() == "markdown":
        return True
    return False


def _adaptive_chunk_window(text: str, base_min: int, base_max: int) -> Tuple[int, int]:
    base_min = max(16, int(base_min))
    base_max = max(base_min + 16, int(base_max))
    approx_tokens = len(_TOKEN_REGEX.findall(text)) or max(1, len(text) // 4)

    if approx_tokens <= base_max:
        min_tokens = max(16, int(base_min * 0.5))
        max_tokens = max(min_tokens + 24, int(base_max * 0.75))
        return min_tokens, max_tokens

    if approx_tokens <= base_max * 3:
        return base_min, base_max

    scale = min(2.0, approx_tokens / float(base_max * 3))
    min_tokens = int(base_min * (1.0 + (scale * 0.5)))
    max_tokens = int(base_max * (1.0 + (scale * 0.5)))
    min_tokens = max(base_min, min(min_tokens, 320))
    max_tokens = max(min_tokens + 32, min(1200, max_tokens))
    remainder = approx_tokens - max_tokens
    if remainder and remainder < min_tokens:
        adjustment = min_tokens - remainder
        max_tokens = max(min_tokens + 32, max_tokens - adjustment)
    return min_tokens, max_tokens


def _apply_uniform_chunks(
    df: "pd.DataFrame",
    *,
    min_tokens: int = DEFAULT_CHUNK_MIN_TOKENS,
    max_tokens: int = DEFAULT_CHUNK_MAX_TOKENS,
    overlap_tokens: int = 0,
) -> "pd.DataFrame":
    if pd is None or df is None or df.empty or "text" not in df.columns:
        return df

    records = df.to_dict(orient="records")
    chunked: List[Dict[str, Any]] = []
    chunker = SemanticChunker(
        max_tokens=max_tokens, 
        overlap_tokens=overlap_tokens,
        min_tokens=min_tokens
    )

    for record in records:
        base_text = str(record.get("text") or "")
        original_text = record.get("text_original") or ""
        doc_hash = record.get("doc_hash", "")
        
        # Use Semantic Chunker
        chunks = chunker.chunk_text(base_text)
        
        if not chunks:
            # Handle empty case
            new_rec = dict(record)
            new_rec["chunk_id"] = 1
            new_rec["chunk_count"] = 1
            new_rec["chunk_tokens"] = 0
            new_rec["text"] = _remove_stopwords(base_text)
            new_rec["content_hash"] = _hash_text(new_rec["text"])
            chunked.append(new_rec)
            continue
            
        chunk_count = len(chunks)
        for i, ch in enumerate(chunks, start=1):
            new_rec = dict(record)
            new_rec["chunk_id"] = i
            new_rec["chunk_count"] = chunk_count
            new_rec["chunk_tokens"] = ch.token_count
            
            # Text cleaning for model
            cleaned_text = ch.text.strip()
            # Restore heading if present
            if "heading" in ch.metadata:
                new_rec["heading"] = ch.metadata["heading"]
                # Optional: prepend heading to text for better embedding context
                # cleaned_text = f"{ch.metadata['heading']}: {cleaned_text}"
            
            new_rec["text"] = cleaned_text
            
            # Map back to original text slice if possible
            # Simplified: Use chunk text as original since we don't have exact char mapping easily
            # logic in SemanticChunker does track char offsets but legacy logic used _slice_text_by_ratio
            if ch.start_char is not None and ch.end_char is not None and original_text:
                # Naive slice mapping if lengths match, else just use chunk text
                # Ideally SemanticChunker should handle original text mapping too
                pass

            new_rec["text_original"] = ch.text 
            new_rec["preview"] = ch.text[:360]
            new_rec["doc_hash"] = doc_hash
            new_rec["content_hash"] = _hash_text(cleaned_text)
            
            chunked.append(new_rec)

    chunks_df = pd.DataFrame(chunked)
    if "target_embed_dtype" in df.attrs:
        chunks_df.attrs["target_embed_dtype"] = df.attrs["target_embed_dtype"]
    return chunks_df


def _time_tokens(epoch: Optional[float]) -> List[str]:
    if not epoch:
        return []
    try:
        dt = datetime.fromtimestamp(float(epoch))
    except Exception:
        return []
    parts = [
        dt.strftime("%Y"),
        dt.strftime("%Y-%m"),
        dt.strftime("%Y-%m-%d"),
        dt.strftime("%B"),
        dt.strftime("%m"),
    ]
    return parts


def _size_bucket(size: Optional[int]) -> Optional[str]:
    if size is None:
        return None
    try:
        size = int(size)
    except (TypeError, ValueError):
        return None
    if size <= 0:
        return None
    if size < 10 * 1024:
        return "size:tiny"
    if size < 1 * 1024 * 1024:
        return "size:small"
    if size < 10 * 1024 * 1024:
        return "size:medium"
    if size < 50 * 1024 * 1024:
        return "size:large"
    return "size:huge"


def _metadata_text(
    path: str,
    ext: str,
    drive: str,
    size: Optional[int] = None,
    mtime: Optional[float] = None,
    ctime: Optional[float] = None,
    owner: Optional[str] = None,
    extra: Optional[str] = None,
) -> str:
    tokens: List[str] = []
    extra_clean: Optional[str] = None
    if path:
        try:
            p = Path(path)
        except Exception:
            p = None
        if p:
            name = p.name
            if name:
                tokens.append(name)
            stem = p.stem
            if stem and stem != name:
                tokens.append(stem)
        else:
            tokens.append(str(path))
    if ext:
        ext_clean = str(ext).strip()
        if ext_clean:
            tokens.append(ext_clean)
            ext_no_dot = ext_clean.lstrip(".")
            if ext_no_dot:
                tokens.append(ext_no_dot)
    if drive:
        drive_str = str(drive)
        tokens.append(drive_str)
    for epoch in (mtime, ctime):
        tokens.extend(_time_tokens(epoch))
    bucket = _size_bucket(size)
    if bucket:
        tokens.append(bucket)
    if owner:
        tokens.append(str(owner))
    if extra:
        extra_clean = TextCleaner.clean(str(extra))
        if extra_clean:
            tokens.extend(_split_tokens(extra_clean))

    seen = set()
    normalized: List[str] = []
    for token in tokens:
        cleaned = TextCleaner.clean(str(token)).lower()
        if not cleaned:
            continue
        if cleaned not in seen:
            seen.add(cleaned)
            normalized.append(cleaned)
    metadata_text = " ".join(normalized)
    if extra_clean:
        return f"{metadata_text}\n{extra_clean}" if metadata_text else extra_clean
    return metadata_text


def _compose_model_text(base_text: str, metadata: str) -> str:
    base = (base_text or "").strip()
    meta = (metadata or "").strip()
    if base:
        if meta and len(base) < 40:
            return f"{base}\n\n{meta}"
        return base
    return meta


_DOC_TAG_HINTS: Tuple[Tuple[str, Tuple[str, ...], Tuple[str, ...]], ...] = (
    (
        "law",
        (
            "법령",
            "법률",
            "법규",
            "조례",
            "규정",
            "규칙",
            "지침",
            "세칙",
            "훈령",
            "행정규칙",
            "regulation",
            "ordinance",
            "bylaw",
        ),
        ("법령문서", "법규", "조례"),
    ),
    (
        "notice",
        (
            "공고",
            "공지",
            "고시",
            "입찰",
            "계약",
            "발주",
            "제안요청서",
            "rfp",
            "announcement",
            "notice",
        ),
        ("공고문", "입찰", "공지문"),
    ),
    (
        "report",
        (
            "보고서",
            "분석",
            "결과보고",
            "백서",
            "리포트",
            "report",
            "analysis",
        ),
        ("보고서", "분석자료"),
    ),
    (
        "minutes",
        (
            "회의록",
            "의사록",
            "minutes",
            "meeting minutes",
        ),
        ("회의록", "회의기록"),
    ),
    (
        "plan",
        (
            "계획",
            "전략",
            "로드맵",
            "마스터플랜",
            "plan",
            "strategy",
            "roadmap",
        ),
        ("계획서", "전략문서"),
    ),
    (
        "manual",
        (
            "매뉴얼",
            "지침서",
            "가이드",
            "guide",
            "manual",
        ),
        ("지침서", "가이드"),
    ),
)


def _infer_doc_tags(path_value: str, extra: Optional[str]) -> Tuple[List[str], List[str]]:
    haystack_parts = []
    if path_value:
        haystack_parts.append(str(path_value))
        try:
            path_obj = Path(path_value)
            haystack_parts.append(path_obj.name)
            haystack_parts.append(path_obj.stem)
        except Exception:
            pass
    if extra:
        haystack_parts.append(str(extra))
    haystack = " ".join(part for part in haystack_parts if part).lower()
    if not haystack:
        return [], []
    tags: List[str] = []
    tokens: List[str] = []
    for slug, keywords, tag_tokens in _DOC_TAG_HINTS:
        for keyword in keywords:
            if keyword.lower() in haystack:
                tags.append(slug)
                for token in tag_tokens:
                    cleaned = TextCleaner.clean(token)
                    if cleaned:
                        tokens.append(cleaned)
                break
    if not tags:
        return [], []
    # deduplicate while preserving order
    seen_tags = set()
    ordered_tags: List[str] = []
    for tag in tags:
        if tag not in seen_tags:
            seen_tags.add(tag)
            ordered_tags.append(tag)
    seen_tokens = set()
    ordered_tokens: List[str] = []
    for token in tokens:
        if token not in seen_tokens:
            seen_tokens.add(token)
            ordered_tokens.append(token)
    return ordered_tags, ordered_tokens


def _prepare_text_frame(df: "pd.DataFrame") -> "pd.DataFrame":
    if pd is None or df is None:
        return df
    if df.empty:
        if MODEL_TEXT_COLUMN not in df.columns:
            df[MODEL_TEXT_COLUMN] = pd.Series(dtype=str)
        return df

    for column in ("text", "text_original"):
        if column in df.columns:
            df[column] = df[column].fillna("").astype(str)

    if "text" not in df.columns:
        df["text"] = ""

    paths = df.get("path")
    if paths is None:
        paths = pd.Series([""] * len(df))
    else:
        paths = paths.fillna("").astype(str)

    exts = df.get("ext")
    if exts is None:
        exts = pd.Series([""] * len(df))
    else:
        exts = exts.fillna("").astype(str)

    drives = df.get("drive")
    if drives is None:
        drives = pd.Series([""] * len(df))
    else:
        drives = drives.fillna("").astype(str)

    sizes = df.get("size")
    if sizes is None:
        sizes = pd.Series([0] * len(df))
    else:
        sizes = sizes.fillna(0)

    mtimes = df.get("mtime")
    if mtimes is None:
        mtimes = pd.Series([0.0] * len(df))
    else:
        mtimes = mtimes.fillna(0.0)

    ctimes = df.get("ctime")
    if ctimes is None:
        ctimes = pd.Series([0.0] * len(df))
    else:
        ctimes = ctimes.fillna(0.0)

    owners = df.get("owner")
    if owners is None:
        owners = pd.Series([""] * len(df))
    else:
        owners = owners.fillna("").astype(str)

    base_texts = df["text"].tolist()
    extra_texts = [
        get_metadata_for_path(str(paths.iat[idx]))
        for idx in range(len(df))
    ]
    metadata_list: List[str] = []
    doc_tags: List[List[str]] = []
    doc_primary_tags: List[str] = []
    for idx in range(len(df)):
        tags, tag_tokens = _infer_doc_tags(paths.iat[idx], extra_texts[idx])
        doc_tags.append(tags)
        doc_primary_tags.append(tags[0] if tags else "")
        metadata_value = _metadata_text(
            paths.iat[idx],
            exts.iat[idx],
            drives.iat[idx],
            size=sizes.iat[idx],
            mtime=mtimes.iat[idx],
            ctime=ctimes.iat[idx],
            owner=owners.iat[idx],
            extra=extra_texts[idx],
        )
        if tag_tokens:
            tag_text = " ".join(tag_tokens)
            metadata_value = f"{metadata_value} {tag_text}".strip() if metadata_value else tag_text
        metadata_list.append(metadata_value)
    df["doc_tags"] = doc_tags
    df["doc_primary_tag"] = doc_primary_tags
    df[MODEL_TEXT_COLUMN] = [
        _compose_model_text(base_texts[idx], metadata_list[idx])
        for idx in range(len(df))
    ]
    return df


def _deduplicate_corpus(df: "pd.DataFrame") -> "pd.DataFrame":
    """Remove duplicate chunks based on content hash (and chunk_id when available)."""
    if pd is None or df is None or df.empty:
        return df
    if "content_hash" not in df.columns or "path" not in df.columns:
        return df

    working = df.copy()
    working["_dedup_pref"] = working["path"].apply(
        lambda p: 1 if str(p or "").startswith("/Volumes/") else 0
    )
    working["_dedup_len"] = working["path"].apply(lambda p: len(str(p or "")))
    working["_dedup_order"] = range(len(working))

    subset_cols = ["content_hash"]
    if "chunk_id" in working.columns:
        subset_cols.append("chunk_id")

    working = working.sort_values(subset_cols + ["_dedup_pref", "_dedup_len", "_dedup_order"])
    working = working.drop_duplicates(subset=subset_cols, keep="first")
    working = working.sort_values("_dedup_order").drop(columns=["_dedup_pref", "_dedup_len", "_dedup_order"])
    working = working.reset_index(drop=True)
    return working


def _resolve_kmeans_n_init() -> Union[str, int]:
    """Return MiniBatchKMeans n_init compatible with installed scikit-learn."""
    try:
        parts = (sklearn_version or "0").split(".")
        major = int(parts[0])
        minor = int(parts[1]) if len(parts) > 1 else 0
        if (major, minor) >= (1, 4):
            return "auto"
    except Exception:
        pass
    return 3


# =========================
# 파이프라인 실행 (메인 함수)
# =========================

def run_step2(
    file_rows: List[Dict[str, Any]],
    out_corpus: Path = Path("./corpus.parquet"),
    out_model: Path = Path("./topic_model.joblib"),
    cfg: TrainConfig = TrainConfig(),
    use_tqdm: bool = True,
    translate: bool = False,
    *,
    scan_state_path: Optional[Path] = None,
    chunk_cache_path: Optional[Path] = None,
    skip_extract: bool = False,
    train_embeddings: bool = True,
):
    global tqdm
    original_tqdm = tqdm
    if not use_tqdm:
        tqdm = None

    if cfg is None:
        cfg = TrainConfig()
    else:
        cfg = replace(cfg)

    target_embed_dtype = _resolve_embed_dtype(cfg)
    cfg.embedding_dtype = target_embed_dtype

    chunk_cache = _create_chunk_cache(chunk_cache_path) if chunk_cache_path else None
    scan_state = load_scan_state(scan_state_path) if scan_state_path else None

    try:
        print("=== Step 2 시작: 내용 추출 & 학습 === (번역: " + ("활성" if translate else "비활성") + ")", flush=True)
        t_all = time.time()
        if pd is None:
            raise RuntimeError("pandas 필요")

        total_count = len(file_rows)
        cached_by_state = 0
        force_paths: Optional[Set[str]] = None

        existing_df = _load_existing_corpus(out_corpus)

        if skip_extract:
            if existing_df is None or existing_df.empty:
                raise RuntimeError(
                    "skip_extract가 설정되었지만 기존 corpus가 없습니다. 먼저 추출을 포함한 pipeline/train을 실행해 corpus.parquet을 생성하세요."
                )
            print("⏭️ 추출 스킵: 기존 corpus를 그대로 사용합니다.", flush=True)
            df = existing_df.copy()
            _prepare_text_frame(df)
            to_process = []
            reused_df = None
            df_new = df.copy()
            df_new_chunks = df.copy()
            process_count = 0
        else:
            if scan_state_path and scan_state is not None:
                forced_rows, cached_rows = filter_incremental_rows(file_rows, scan_state)
                force_paths = {str(row.get("path") or "") for row in forced_rows if row.get("path")}
                cached_by_state = len(cached_rows)
                cached_paths = {str(row.get("path") or "") for row in cached_rows if row.get("path")}
                if force_paths:
                    print(
                        f"⚙️ 증분 상태: {len(force_paths):,}건 재처리, 캐시 일치 {cached_by_state:,}건",
                        flush=True,
                    )
                else:
                    print("⚙️ 증분 상태: 신규 변경 없음", flush=True)

            to_process, reused_df = _split_cache(file_rows, existing_df, force_paths=force_paths)
            if scan_state_path and scan_state is not None:
                # If scan-state marks a path as cached, never re-extract it even if the corpus
                # no longer contains that path (e.g. chunk-level deduplication kept only one copy).
                to_process = [
                    row
                    for row in to_process
                    if str(row.get("path") or "") and str(row.get("path") or "") not in cached_paths
                ]
            process_paths = {str(row.get("path") or "") for row in to_process if row.get("path")}
            process_count = len(process_paths)
            print(
                f"🗃️ 신규/변경 추출 대상: {process_count:,} | 총 스캔: {total_count:,}",
                flush=True,
            )

            if process_count == 0:
                if reused_df is not None:
                    df = reused_df.copy()
                elif existing_df is not None:
                    df = existing_df.copy()
                else:
                    df = pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))
                _prepare_text_frame(df)
                order_map = {row["path"]: idx for idx, row in enumerate(file_rows)} if file_rows else {}
                if "path" in df.columns and order_map:
                    df["_order"] = df["path"].map(order_map)
                    df = df.sort_values("_order").drop(columns=["_order"]).reset_index(drop=True)
                df = _deduplicate_corpus(df)
                CorpusBuilder.save(df, out_corpus)
                if chunk_cache:
                    chunk_cache.update_from_frame(df)
                    chunk_cache.save()
                if scan_state_path:
                    updated_state = update_scan_state(scan_state or {}, file_rows)
                    save_scan_state(scan_state_path, updated_state)
                df.attrs["metrics"] = {}
                df.attrs["incremental"] = {
                    "requested": process_count,
                    "effective": 0,
                    "skipped_by_state": cached_by_state,
                    "total": total_count,
                }
                df.attrs["target_embed_dtype"] = target_embed_dtype
                print("✨ 변경된 문서가 없어 기존 모델을 유지합니다.", flush=True)
                return df, None

            cb = CorpusBuilder(
                max_text_chars=200_000,
                progress=use_tqdm,
                translate=translate,
                target_embed_dtype=target_embed_dtype,
                # PyMuPDF가 다중 스레드에서 불안정하므로 macOS 기본은 워커 1개로 제한
                max_workers=int(os.getenv("INFOPILOT_MAX_EXTRACT_WORKERS", "1")),
            )
            df_new = cb.build(to_process) if process_count else pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

        restored_df = None
        unchanged_paths: Set[str] = set()
        if chunk_cache and df_new is not None and not df_new.empty:
            unchanged_paths = chunk_cache.unchanged_paths(df_new)
            if unchanged_paths:
                print(f"♻️ 내용 해시 동일 문서 재사용: {len(unchanged_paths):,}", flush=True)
                if existing_df is not None:
                    restored_df = _collect_existing_rows(existing_df, unchanged_paths)
                df_new = df_new[~df_new["path"].isin(list(unchanged_paths))]

            df_new_chunks = (
                _apply_uniform_chunks(
                    df_new,
                    min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
                    max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
                )
                if df_new is not None and not df_new.empty
                else pd.DataFrame(columns=list(df_new.columns) if df_new is not None else list(ExtractRecord.__annotations__.keys()))
            )
            if hasattr(df_new_chunks, "attrs"):
                df_new_chunks.attrs["target_embed_dtype"] = target_embed_dtype

            frames: List["pd.DataFrame"] = []
            if reused_df is not None and not reused_df.empty:
                frames.append(reused_df)
            if restored_df is not None and not restored_df.empty:
                frames.append(restored_df)
            if df_new_chunks is not None and not df_new_chunks.empty:
                frames.append(df_new_chunks)

            if frames:
                df = pd.concat(frames, ignore_index=True)
            else:
                df = pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

            _prepare_text_frame(df)

            order_map = {row["path"]: idx for idx, row in enumerate(file_rows)} if file_rows else {}
            if "path" in df.columns and order_map:
                df["_order"] = df["path"].map(order_map)
                df = df.sort_values("_order").drop(columns=["_order"]).reset_index(drop=True)

            if "ok" in df.columns:
                df["ok"] = df["ok"].apply(lambda v: bool(v) if isinstance(v, bool) else str(v).strip().lower() in {"true", "1", "yes"})
            if "topic" in df.columns:
                df = df.drop(columns=["topic"])

        text_col = MODEL_TEXT_COLUMN if MODEL_TEXT_COLUMN in df.columns else "text"
        text_mask = df[text_col].fillna("").str.len() > 0
        train_df = df[df["ok"] & text_mask].copy()
        if not train_df.empty:
            _prepare_text_frame(train_df)
        print(f"🧹 학습 대상 문서: {len(train_df):,}/{len(df):,}", flush=True)
        if len(train_df) == 0:
            df = _deduplicate_corpus(df)
            CorpusBuilder.save(df, out_corpus)
            if scan_state_path:
                updated_state = update_scan_state(scan_state or {}, file_rows)
                save_scan_state(scan_state_path, updated_state)
            if chunk_cache:
                chunk_cache.update_from_frame(df)
                chunk_cache.save()
            print(f"⚠️ 유효 텍스트 없음. 코퍼스만 저장: {out_corpus}", flush=True)
            df.attrs["metrics"] = {}
            df.attrs["incremental"] = {
                "requested": process_count,
                "effective": 0,
                "skipped_by_state": cached_by_state,
                "total": total_count,
            }
            df.attrs["target_embed_dtype"] = target_embed_dtype
            return df, None

        if not train_embeddings:
            df = _deduplicate_corpus(df)
            CorpusBuilder.save(df, out_corpus)
            if scan_state_path:
                updated_state = update_scan_state(scan_state or {}, file_rows)
                save_scan_state(scan_state_path, updated_state)
            if chunk_cache:
                chunk_cache.update_from_frame(df)
                chunk_cache.save()
            print(f"📦 추출만 완료 (임베딩/모델 건너뜀): {out_corpus}", flush=True)
            df.attrs["metrics"] = {}
            df.attrs["incremental"] = {
                "requested": process_count,
                "effective": 0,
                "skipped_by_state": cached_by_state,
                "total": total_count,
            }
            df.attrs["target_embed_dtype"] = target_embed_dtype
            return df, None

        topics_df = None
        model_obj: Optional[Any] = None
        metrics: Dict[str, float] = {}

        if cfg.use_sentence_transformer and SentenceTransformer is not None:
            try:
                if cfg.embedding_chunk_size and cfg.embedding_chunk_size > 0:
                    embeddings, semantic_model, metrics = _fit_sentence_transformer_chunked(
                        train_df, text_col, cfg
                    )
                else:
                    semantic_model = SentenceBertModel(cfg)
                    embeddings = semantic_model.fit(train_df, text_col=text_col)
                    if semantic_model.cluster_labels_ is not None:
                        metrics = evaluate_embeddings(
                            embeddings,
                            semantic_model.cluster_labels_,
                            topk=min(5, max(1, embeddings.shape[0] - 1)),
                        )
                print(
                    f"✅ Sentence-BERT 임베딩 완료 (docs={embeddings.shape[0]:,}, dim={semantic_model.embedding_dim})",
                    flush=True,
                )
                if semantic_model.cluster_labels_ is not None:
                    train_df["topic"] = semantic_model.cluster_labels_
                    topics_df = train_df[["path", "topic"]].copy()
                model_obj = semantic_model
            except Exception as exc:
                raise RuntimeError(
                    f"Sentence-BERT 임베딩에 실패했습니다. TF-IDF 백업을 사용하지 않습니다: {exc}"
                ) from exc
        elif cfg.use_sentence_transformer and SentenceTransformer is None:
            raise RuntimeError("sentence-transformers가 설치되어 있지 않아 임베딩을 진행할 수 없습니다.")

        if model_obj is None:
            tm = TopicModel(cfg)
            tm.fit(train_df, text_col=text_col)
            labels = tm.predict(train_df, text_col=text_col)
            train_df["topic"] = labels
            topics_df = train_df[["path", "topic"]].copy()
            model_obj = tm
            metrics = {}

        if topics_df is not None:
            df = df.merge(topics_df, on="path", how="left")

        df = _deduplicate_corpus(df)
        CorpusBuilder.save(df, out_corpus)

        if isinstance(model_obj, SentenceBertModel):
            model_obj.save(out_model)
        elif isinstance(model_obj, TopicModel) and joblib:
            model_obj.save(out_model)

        if chunk_cache:
            current_paths = set(df["path"].astype(str)) if "path" in df.columns else set()
            missing = chunk_cache.known_paths() - current_paths
            if missing:
                chunk_cache.drop_paths(missing)
            chunk_cache.update_from_frame(df)
            chunk_cache.save()

        if scan_state_path:
            updated_state = update_scan_state(scan_state or {}, file_rows)
            save_scan_state(scan_state_path, updated_state)

        dt_all = time.time() - t_all
        print(f"💾 저장 완료: corpus → {out_corpus} | model → {out_model}", flush=True)
        print(f"🎉 Step 2 종료 (총 {dt_all:.1f}s)", flush=True)
        df.attrs["metrics"] = metrics or {}
        df.attrs["incremental"] = {
            "requested": process_count,
            "effective": len(df_new_chunks["path"].unique()) if not df_new_chunks.empty else 0,
            "skipped_by_state": cached_by_state,
            "total": total_count,
        }
        df.attrs["target_embed_dtype"] = target_embed_dtype
        return df, model_obj
    finally:
        tqdm = original_tqdm



def update_corpus_file(
    new_records: "pd.DataFrame",
    corpus_path: Path,
) -> "pd.DataFrame":
    """Merge `new_records` into the persisted corpus and return the updated frame."""
    if pd is None:
        raise RuntimeError("pandas 필요. pip install pandas")

    existing = _load_existing_corpus(corpus_path)
    if existing is None or existing.empty:
        combined = new_records.copy()
    else:
        if "path" in existing.columns and "path" in new_records.columns:
            paths_to_replace = set(new_records["path"].astype(str).tolist())
            mask = ~existing["path"].astype(str).isin(paths_to_replace)
            combined = pd.concat([existing[mask], new_records], ignore_index=True)
        else:
            combined = pd.concat([existing, new_records], ignore_index=True)

    combined = _apply_uniform_chunks(
        combined,
        min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
        max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
    )
    _prepare_text_frame(combined)
    combined = _deduplicate_corpus(combined)
    CorpusBuilder.save(combined, corpus_path)
    return combined


def remove_from_corpus(paths: List[str], corpus_path: Path) -> "pd.DataFrame":
    """Remove documents whose paths match `paths` from the persisted corpus."""
    if pd is None:
        raise RuntimeError("pandas 필요. pip install pandas")

    existing = _load_existing_corpus(corpus_path)
    if existing is None or existing.empty:
        return pd.DataFrame(columns=list(ExtractRecord.__annotations__.keys()))

    to_drop = {str(p) for p in paths}
    if "path" not in existing.columns:
        return existing

    filtered = existing[~existing["path"].astype(str).isin(to_drop)].copy()
    filtered = _apply_uniform_chunks(
        filtered,
        min_tokens=DEFAULT_CHUNK_MIN_TOKENS,
        max_tokens=DEFAULT_CHUNK_MAX_TOKENS,
    )
    _prepare_text_frame(filtered)
    filtered = _deduplicate_corpus(filtered)
    CorpusBuilder.save(filtered, corpus_path)
    return filtered
