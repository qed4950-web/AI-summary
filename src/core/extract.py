"""extract module split from pipeline (auto-split from originals)."""
from __future__ import annotations
import re
import sys
from pathlib import Path
from typing import Tuple, Dict, Any
import io
from PIL import Image

# -- Optional Dependencies & Setup --

try:
    import pypandoc
    print("[INFO] pypandoc library found.")
    # --- Auto-install Pandoc if not found ---
    try:
        pandoc_path = pypandoc.get_pandoc_path()
        print(f"[INFO] Pandoc executable found at: {pandoc_path}")
    except OSError:
        print("[WARNING] Pandoc executable not found. Attempting to download...")
        try:
            pypandoc.download_pandoc()
            print("[SUCCESS] Pandoc downloaded and installed successfully.")
        except Exception as e:
            print(f"[ERROR] Failed to download Pandoc: {e}")
    # ----------------------------------------
except ImportError:
    print("[ERROR] pypandoc library not installed.")
    pypandoc = None

try:
    import docx
except ImportError:
    docx = None

try:
    import fitz  # PyMuPDF
except ImportError:
    fitz = None

try:
    import easyocr
except ImportError:
    easyocr = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import pyhwp
except ImportError:
    pyhwp = None

try:
    import openpyxl
    from openpyxl.utils.exceptions import InvalidFileException
except ImportError:
    openpyxl = None
    InvalidFileException = type('InvalidFileException', (Exception,), {})

try:
    import xlrd
except ImportError:
    xlrd = None


class TextCleaner:
    _multi_space = re.compile(r"\s+")
    @classmethod
    def clean(cls, s:str)->str:
        if not s: return ""
        s = s.replace("\x00"," ")
        s = cls._multi_space.sub(" ", s)
        return s.strip()

class BaseExtractor:
    exts: Tuple[str,...] = ()
    def can_handle(self, p:Path)->bool: return p.suffix.lower() in self.exts
    def extract(self, p:Path)->Dict[str,Any]: raise NotImplementedError

class HwpExtractor(BaseExtractor):
    exts=(".hwp",)
    def extract(self, p:Path)->Dict[str,Any]:
        current_errors = []
        if pyhwp:
            try:
                with open(p, "rb") as f:
                    hwp_file = pyhwp.HwpFile(f.read())
                    hwp_file.parse_header()
                    hwp_file.parse_body()
                    text = pyhwp.utils.get_text(hwp_file)
                    if text.strip():
                        return {"ok":True,"text":TextCleaner.clean(text)}
                    else:
                        current_errors.append("pyhwp HWP 추출했으나 내용이 비어있음")
            except Exception as e:
                current_errors.append(f"pyhwp HWP 추출 실패: {e}")
        else:
            current_errors.append("pyhwp 라이브러리 미설치")

        return {"ok":False,"text":f"HWP 추출 실패 ({'; '.join(current_errors)})"}

class DocxExtractor(BaseExtractor):
    exts=(".docx",".doc")
    def extract(self, p:Path)->Dict[str,Any]:
        current_errors = []

        if p.suffix.lower() == ".docx":
            if docx:
                try:
                    d = docx.Document(str(p)); t = "\n".join(par.text for par in d.paragraphs)
                    if t.strip(): return {"ok":True,"text":TextCleaner.clean(t)}
                    else: current_errors.append("python-docx로 DOCX 추출했으나 내용이 비어있음")
                except Exception as e: current_errors.append(f"python-docx DOCX 추출 실패: {e}")
            else: current_errors.append("'python-docx' 라이브러리 미설치")

        elif p.suffix.lower() == ".doc":
            if pypandoc:
                try:
                    text = pypandoc.convert_file(str(p), 'plain')
                    if text.strip(): return {"ok": True, "text": TextCleaner.clean(text)}
                    else: current_errors.append("pypandoc으로 DOC 추출했으나 내용이 비어있음")
                except Exception as e: current_errors.append(f"pypandoc DOC 추출 실패: {e}")
            else:
                current_errors.append("pypandoc 라이브러리 미설치 (DOC 파일 처리에 필요)")
        
        return {"ok":False,"text":f"DOC/DOCX 추출 실패 ({'; '.join(current_errors)})"}

class PdfExtractor(BaseExtractor):
    exts=(".pdf",)
    _ocr_reader = None  # Singleton for EasyOCR reader

    @classmethod
    def get_ocr_reader(cls):
        """Initializes and returns a singleton EasyOCR reader instance."""
        if cls._ocr_reader is None:
            if easyocr:
                print("Initializing EasyOCR Reader...")
                cls._ocr_reader = easyocr.Reader(['ko', 'en'], gpu=True)
        return cls._ocr_reader

    def extract(self, p:Path)->Dict[str,Any]:
        current_errors = []
        text_content = ""
        doc = None

        # 1. Attempt text extraction with PyMuPDF
        if fitz:
            try:
                doc = fitz.open(p)
                if doc.is_encrypted:
                    current_errors.append("PDF가 암호화되어 있습니다.")
                else:
                    text_content = "\n".join(page.get_text() for page in doc)
                    if text_content.strip():
                        doc.close()
                        return {"ok":True, "text":TextCleaner.clean(text_content)}
                    else:
                        current_errors.append("PyMuPDF로 추출했으나 내용이 비어있음")
            except Exception as e:
                current_errors.append(f"PyMuPDF PDF 추출 실패: {e}")
                if doc: doc.close()
                doc = None
        else:
            current_errors.append("PyMuPDF(fitz) 라이브러리 미설치")

        # 2. Fallback to EasyOCR if PyMuPDF found no text
        if not text_content.strip():
            if easyocr:
                reader = self.get_ocr_reader()
                if not reader:
                    current_errors.append("EasyOCR 리더기 초기화 실패")
                else:
                    try:
                        if doc is None: doc = fitz.open(p)
                        ocr_text_parts = []
                        for page in doc:
                            pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
                            img_bytes = pix.tobytes("png")
                            result = reader.readtext(img_bytes, detail=0, paragraph=True)
                            ocr_text_parts.extend(result)
                        text_content = "\n".join(ocr_text_parts)
                        if text_content.strip():
                            return {"ok":True, "text":TextCleaner.clean(text_content)}
                        else:
                            current_errors.append("EasyOCR 추출했으나 내용이 비어있음")
                    except Exception as e:
                        current_errors.append(f"EasyOCR PDF 추출 실패: {e}")
            else:
                current_errors.append("easyocr 라이브러리 미설치")
        
        if doc: doc.close()

        final_error_msg = "PDF 추출 실패"
        real_errors = [e for e in current_errors if "내용이 비어있음" not in e]
        if real_errors:
             final_error_msg += f" ({'; '.join(real_errors)})"
        elif current_errors:
             final_error_msg += f" ({current_errors[-1]})"

        return {"ok":False,"text": final_error_msg}

class PptxExtractor(BaseExtractor):
    exts=(".pptx",".ppt")
    def extract(self, p:Path)->Dict[str,Any]:
        current_errors = []

        if p.suffix.lower() == ".pptx":
            if pptx:
                try:
                    pres = pptx.Presentation(str(p))
                    texts = [shape.text for slide in pres.slides for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
                    if texts: return {"ok":True,"text":TextCleaner.clean("\n".join(texts))}
                    else: current_errors.append("python-pptx로 PPTX 추출했으나 내용이 비어있음")
                except Exception as e: current_errors.append(f"python-pptx PPTX 추출 실패: {e}")
            else: current_errors.append("'python-pptx' 라이브러리 미설치")

        elif p.suffix.lower() == ".ppt":
            if pypandoc:
                try:
                    text = pypandoc.convert_file(str(p), 'plain')
                    if text.strip(): return {"ok": True, "text": TextCleaner.clean(text)}
                    else: current_errors.append("pypandoc으로 PPT 추출했으나 내용이 비어있음")
                except Exception as e: current_errors.append(f"pypandoc PPT 추출 실패: {e}")
            else:
                current_errors.append("pypandoc 라이브러리 미설치 (PPT 파일 처리에 필요)")

        return {"ok":False,"text":f"PPT/PPTX 추출 실패 ({'; '.join(current_errors)})"}

class ExcelExtractor(BaseExtractor):
    exts = (".xlsx", ".xlsm", ".xls")
    def extract(self, p: Path) -> Dict[str, Any]:
        current_errors = []
        text_content = []

        try:
            if p.suffix.lower() in (".xlsx", ".xlsm"):
                if not openpyxl: raise ImportError("openpyxl 라이브러리 미설치")
                workbook = openpyxl.load_workbook(p, data_only=True)
                for sheet in workbook.worksheets:
                    for row in sheet.iter_rows():
                        text_content.extend([str(cell.value) for cell in row if cell.value is not None])
                if not text_content: current_errors.append("openpyxl로 추출했으나 내용이 비어있음")
            elif p.suffix.lower() == ".xls":
                if not xlrd: raise ImportError("xlrd 라이브러리 미설치")
                workbook = xlrd.open_workbook(p)
                for sheet in workbook.sheets():
                    for row_idx in range(sheet.nrows):
                        text_content.extend([str(sheet.cell_value(row_idx, col_idx)) for col_idx in range(sheet.ncols) if sheet.cell_value(row_idx, col_idx) is not None])
                if not text_content: current_errors.append("xlrd로 추출했으나 내용이 비어있음")
            else: current_errors.append(f"지원되지 않는 Excel 확장자: {p.suffix.lower()}")
        except Exception as e: current_errors.append(f"Excel 추출 실패: {e}")
        
        if text_content and not current_errors:
            return {"ok": True, "text": TextCleaner.clean(" ".join(text_content))}
        else:
            return {"ok": False, "text": f"Excel 추출 실패 ({'; '.join(current_errors)})"}

class TextExtractor(BaseExtractor):
    exts = (".csv", ".txt")
    def extract(self, p: Path) -> Dict[str, Any]:
        tried_encodings = ["utf-8", "cp949", "latin-1"]
        for encoding in tried_encodings:
            try:
                with open(p, "r", encoding=encoding) as f:
                    text = f.read()
                return {"ok": True, "text": TextCleaner.clean(text)}
            except (UnicodeDecodeError, PermissionError):
                continue
            except Exception as e:
                return {"ok": False, "text": f"텍스트 파일 추출 실패: {e}"}
        return {"ok": False, "text": "모든 인코딩으로 텍스트 파일 디코딩 실패"}

EXT_MAP = {
    ext: extractor_class()
    for extractor_class in [HwpExtractor, DocxExtractor, PdfExtractor, PptxExtractor, ExcelExtractor, TextExtractor]
    for ext in extractor_class.exts
}